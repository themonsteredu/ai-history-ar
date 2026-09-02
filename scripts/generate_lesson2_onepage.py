from __future__ import annotations

import argparse
import json
import re
import zipfile
from pathlib import Path
from typing import Any, Iterable

from fontTools.ttLib import TTFont as FontToolsTTFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pypdf import PdfReader, PdfWriter
from reportlab.lib import colors
from reportlab.lib.colors import HexColor
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph, Table, TableStyle

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = ROOT / "public" / "downloads"
DEFAULT_FONTS = ROOT / "public" / "fonts"

PRIMARY = HexColor("#2F715B")
PRIMARY_DARK = HexColor("#1F4F40")
MINT = HexColor("#EAF3EF")
MINT_LIGHT = HexColor("#F6FAF8")
BEIGE = HexColor("#FBF2E2")
GRID = HexColor("#C8D2CE")
TEXT = HexColor("#1F2925")
MUTED = HexColor("#64716C")
WHITE = colors.white

# 학생 활동지는 학교에 제출하는 흑백 양식을 씁니다(1차시 활동지와 같은 팔레트).
FORM_INK = HexColor("#111111")
FORM_HEAD = HexColor("#222222")
FORM_BAND = HexColor("#F0F0F0")
FORM_LINE = HexColor("#B8B8B8")
FORM_SOFT = HexColor("#444444")

PPT_PRIMARY = RGBColor(47, 113, 91)
PPT_DARK = RGBColor(31, 79, 64)
PPT_MINT = RGBColor(234, 243, 239)
PPT_BEIGE = RGBColor(251, 242, 226)
PPT_TEXT = RGBColor(31, 41, 37)
PPT_MUTED = RGBColor(100, 113, 108)

# Each statement: text, verdict, teacher rationale, suggested search words.
DATA: dict[str, dict[str, Any]] = {
    "three-kingdoms": {
        "short": "삼국시대",
        "course": "삼국시대 문화유산 해설사",
        "core": "1500년 전 이야기, AI는 어떻게 알고 있을까?",
        "next": "출처 · 시기 · 교차 · 원본 · 보류",
        "delta": "의견 나뉨·근거 부족",
        "appTab": "활동지 수업",
        "appPath": "/three-kingdoms/lesson/2?view=activity",
        "sourceRanks": ["1순위  국가유산청 국가유산포털", "2순위  국립중앙박물관·국립부여박물관", "3순위  유네스코 세계유산센터·우리역사넷", "비교용  블로그·영상·AI 답변"],
        "groups": [
            {
                "id": 1, "slug": "muryeongwangneung", "heritage": "무령왕릉", "color": "#7A5AA6",
                "focus": "발굴은 왜 아쉬움으로 남았을까?",
                "question": "무령왕릉은 어떻게 발견되었고, 무엇을 알게 되었나요?",
                "source": "국립공주박물관·국립중앙박물관 무령왕릉 자료",
                "items": [
                    ("무령왕릉은 1971년 공주 송산리 고분군의 배수로 공사 중 우연히 발견되었습니다.", "○ 자료로 확인", "발견 연도와 경위가 공식 자료에 기록되어 있습니다.", "무령왕릉 1971 배수로 공사"),
                    ("무덤 안의 지석 덕분에 무덤의 주인이 무령왕과 왕비라는 사실을 확인할 수 있었습니다.", "○ 자료로 확인", "지석에 무덤 주인과 장례 정보가 남아 있습니다.", "무령왕릉 지석 무령왕 왕비"),
                    ("발견 당시 이미 여러 차례 도굴되어 중요한 유물 대부분이 사라진 상태였습니다.", "× 자료와 다름", "도굴되지 않은 상태로 발견되어 많은 유물이 남아 있었습니다.", "무령왕릉 도굴 여부"),
                    ("무령왕릉을 처음 발견한 사람은 현장 학습을 온 초등학생이었습니다.", "× 자료와 다름", "배수로 공사 과정에서 공사 관계자들이 발견했습니다.", "무령왕릉 발견 과정"),
                    ("무덤 안에서는 진묘수, 금제 관식, 금동 신발, 청동 거울 등 여러 유물이 나왔습니다.", "○ 자료로 확인", "대표 출토 유물은 박물관 자료에서 확인됩니다.", "무령왕릉 출토 유물"),
                    ("지석만 읽으면 장례의 모든 과정과 유물 하나하나의 정확한 쓰임까지 완전히 알 수 있습니다.", "△ 의견 나뉨·근거 부족", "지석은 중요한 단서지만 모든 과정과 쓰임을 전부 설명하지는 않습니다.", "무령왕릉 지석 한계"),
                ],
            },
            {
                "id": 2, "slug": "baekje-incense-burner", "heritage": "백제 금동대향로", "color": "#3D7F5D",
                "focus": "향로에 새겨진 것은 무엇을 뜻할까?",
                "question": "백제 금동대향로는 어디에서 발견되었고, 어떤 모습인가요?",
                "source": "국립부여박물관 백제 금동대향로 자료",
                "items": [
                    ("백제 금동대향로는 1993년 부여 능산리 절터를 발굴하던 중 발견되었습니다.", "○ 자료로 확인", "발견 연도와 장소가 박물관 자료에 기록되어 있습니다.", "금동대향로 1993 능산리"),
                    ("이 향로는 청동으로 만든 뒤 표면에 금을 입힌 금동 유물입니다.", "○ 자료로 확인", "금동은 청동 바탕에 금을 입힌 재료를 뜻합니다.", "금동대향로 재료"),
                    ("향로는 백제 왕의 무덤 안에서 발견되었습니다.", "× 자료와 다름", "왕릉이 아니라 능산리 절터에서 출토되었습니다.", "금동대향로 발견 장소"),
                    ("향로의 꼭대기에는 봉황, 받침에는 용이 있고 산·연꽃·동물·악사도 표현되어 있습니다.", "○ 자료로 확인", "향로의 형태와 장식은 소장품 사진과 설명에서 확인됩니다.", "금동대향로 봉황 용 악사"),
                    ("향로에 새겨진 모든 동물과 인물이 무엇을 뜻하는지 당시 기록에 남아 있어 완전히 밝혀졌습니다.", "△ 의견 나뉨·근거 부족", "형상은 관찰되지만 모든 무늬의 뜻을 하나로 확정하기 어렵습니다.", "금동대향로 무늬 의미"),
                    ("이 향로는 신라의 황룡사에서 만들어져 발견된 유물입니다.", "× 자료와 다름", "부여 지역 절터에서 발견된 백제 유물입니다.", "금동대향로 백제 신라"),
                ],
            },
            {
                "id": 3, "slug": "cheomseongdae", "heritage": "첨성대", "color": "#345D8C",
                "focus": "첨성대는 정말 천문대였을까?",
                "question": "첨성대는 무엇을 하던 곳이며, 사용 방법은 모두 밝혀졌나요?",
                "source": "국가유산청 경주 첨성대 안내",
                "items": [
                    ("첨성대는 경주에 남아 있는 돌로 쌓은 신라시대 건축물입니다.", "○ 자료로 확인", "경주에 있는 석조 건축물임을 확인할 수 있습니다.", "경주 첨성대 석조"),
                    ("국가유산 안내에서는 첨성대를 선덕여왕 때 세운 천문 관측대로 설명합니다.", "○ 자료로 확인", "공식 안내에 건립 시기와 천문 관측 관련 설명이 있습니다.", "첨성대 선덕여왕 천문"),
                    ("첨성대 꼭대기에는 별을 확대해서 보는 망원경이 설치되어 있었습니다.", "× 자료와 다름", "망원경은 첨성대가 세워진 시대보다 훨씬 뒤에 등장했습니다.", "첨성대 망원경"),
                    ("첨성대는 나무 기둥과 기와로만 만든 건물이었습니다.", "× 자료와 다름", "현재 남은 첨성대는 돌을 층층이 쌓은 구조물입니다.", "첨성대 재료 돌"),
                    ("첨성대의 돌 개수와 모든 층의 수가 달력을 뜻한다는 사실이 옛 문헌으로 완전히 증명되었습니다.", "△ 의견 나뉨·근거 부족", "숫자의 상징에는 여러 해석이 있어 확정된 사실로 보기 어렵습니다.", "첨성대 돌 개수 학설"),
                    ("신라 사람들이 첨성대 안팎에서 별을 관찰한 정확한 방법은 학자들이 모두 같은 결론을 내렸습니다.", "△ 의견 나뉨·근거 부족", "구체적인 사용 방법에는 여러 해석이 있습니다.", "첨성대 사용 방법 학설"),
                ],
            },
            {
                "id": 4, "slug": "silla-gold-crown", "heritage": "신라 금관", "color": "#B18A2E",
                "focus": "신라 금관은 실제로 쓰고 다녔을까?",
                "question": "신라 금관은 누가, 언제, 어떻게 사용했나요?",
                "source": "국립중앙박물관 신라 금관 자료",
                "items": [
                    ("신라 금관은 왕릉급의 큰 무덤에서 권위를 보여 주는 부장품으로 발견되었습니다.", "○ 자료로 확인", "출토 상황과 금관의 성격은 박물관 자료에서 확인됩니다.", "신라 금관 부장품"),
                    ("금관에는 얇은 금판으로 만든 세움 장식과 굽은옥 같은 장식이 보입니다.", "○ 자료로 확인", "금관의 실제 형태와 장식에서 확인됩니다.", "신라 금관 굽은옥"),
                    ("신라 금관은 오직 왕의 무덤에서만 발견되었으므로 모두 왕이 쓴 왕관입니다.", "× 자료와 다름", "높은 신분의 인물 무덤에서도 발견되어 모두 왕의 관이라 단정할 수 없습니다.", "신라 금관 왕족 무덤"),
                    ("발견된 금관은 모두 무덤 주인의 머리에 똑바로 씌워진 상태였습니다.", "× 자료와 다름", "얼굴이나 가슴 부근 등 서로 다른 위치에서 발견된 사례가 있습니다.", "신라 금관 출토 위치"),
                    ("신라 사람들은 금관을 농사나 일상 노동을 할 때도 매일 쓰고 다녔습니다.", "× 자료와 다름", "일상 작업용으로 사용했다는 근거는 없습니다.", "신라 금관 일상 착용"),
                    ("금관이 생전에 실제로 착용된 것인지 장례를 위해 제작된 것인지에 대해서는 여러 연구 의견이 있습니다.", "△ 의견 나뉨·근거 부족", "실제 착용설과 장례용 해석을 함께 살펴봐야 합니다.", "신라 금관 실제 착용"),
                ],
            },
            {
                "id": 5, "slug": "goguryeo-murals", "heritage": "고구려 고분벽화", "color": "#8A4B3F",
                "focus": "고구려 벽화는 무엇을 그린 것일까?",
                "question": "고구려 고분벽화에는 어떤 장면이 있으며, 무엇을 알 수 있나요?",
                "source": "UNESCO 고구려 고분군·국립박물관 자료",
                "items": [
                    ("고구려 고분벽화에는 사냥, 행렬, 생활 모습, 무늬와 사신도 같은 그림이 남아 있습니다.", "○ 자료로 확인", "벽화의 주제는 공식 자료에서 확인할 수 있습니다.", "고구려 벽화 사냥 사신도"),
                    ("벽화무덤은 주로 왕족과 귀족의 무덤과 관련된 것으로 설명됩니다.", "○ 자료로 확인", "무덤의 규모와 내용을 바탕으로 왕족·귀족 무덤으로 봅니다.", "고구려 벽화무덤 왕족"),
                    ("발견된 고구려 무덤에는 하나도 빠짐없이 모두 벽화가 그려져 있습니다.", "× 자료와 다름", "고구려 무덤 가운데 일부에만 벽화가 남아 있습니다.", "고구려 무덤 모두 벽화"),
                    ("고구려 고분벽화는 모두 현재 대한민국 안에서만 발견되었습니다.", "× 자료와 다름", "북한과 중국 동북 지역에도 고구려 벽화무덤이 있습니다.", "고구려 벽화 북한 중국"),
                    ("벽화에 보이는 장면은 모든 고구려 사람의 하루를 사진처럼 정확히 보여 줍니다.", "△ 의견 나뉨·근거 부족", "왕족·귀족 무덤의 장면을 모든 사람의 생활로 일반화하면 안 됩니다.", "고구려 벽화 생활 일반화"),
                    ("한 명의 유명한 화가가 고구려의 모든 고분벽화를 혼자 그렸다는 기록이 남아 있습니다.", "× 자료와 다름", "모든 벽화를 한 사람이 그렸다는 기록은 확인되지 않습니다.", "고구려 벽화 화가 기록"),
                ],
            },
            {
                "id": 6, "slug": "gaya-tumuli", "heritage": "가야 고분군", "color": "#C87332",
                "focus": "가야는 왜 오랫동안 잊혔을까?",
                "question": "가야 고분군은 어떤 나라의 모습을 보여 주나요?",
                "source": "국가유산청 가야고분군·UNESCO 자료",
                "items": [
                    ("가야 고분군 세계유산은 여러 지역에 있는 일곱 고분군으로 이루어져 있습니다.", "○ 자료로 확인", "세계유산은 일곱 고분군으로 구성되어 있습니다.", "가야 고분군 일곱"),
                    ("가야는 여러 정치체가 함께 존재한 연맹적 성격을 지녔으며 지역마다 차이도 있었습니다.", "○ 자료로 확인", "고분군은 가야 정치체의 공통점과 지역 차이를 보여 줍니다.", "가야 여러 정치체"),
                    ("가야는 처음부터 한 명의 왕이 모든 지역을 다스린 하나의 중앙집권 제국이었습니다.", "× 자료와 다름", "여러 정치체가 공존한 연맹적 성격으로 설명됩니다.", "가야 중앙집권 제국"),
                    ("세계유산에 포함된 가야 고분군 일곱 곳은 모두 같은 한 도시 안에 있습니다.", "× 자료와 다름", "고분군은 여러 지역에 분포합니다.", "가야 고분군 위치"),
                    ("가야 고분군은 2023년에 유네스코 세계유산으로 등재되었습니다.", "○ 자료로 확인", "2023년에 유네스코 세계유산으로 등재되었습니다.", "가야 고분군 2023"),
                    ("고분 하나마다 묻힌 사람의 이름과 일생이 모두 문헌에 남아 있어 정확히 알 수 있습니다.", "△ 의견 나뉨·근거 부족", "문헌이 부족해 무덤 주인과 일생을 모두 특정하기 어렵습니다.", "가야 고분 주인 문헌"),
                ],
            },
        ],
    },
    "joseon": {
        "short": "조선시대",
        "course": "조선시대 문화유산 해설사",
        "core": "우리가 아는 조선, 정말 그랬을까?",
        "next": "출처 · 시기 · 교차 · 원본",
        "delta": "근거 부족·과장됨",
        "appTab": "활동지 수업",
        "appPath": "/joseon/lesson/2?view=activity",
        "sourceRanks": ["1순위  국가유산청 국가유산포털", "2순위  국립고궁박물관·국립중앙박물관", "3순위  국사편찬위원회·조선왕조실록", "비교용  블로그·영상·AI 답변"],
        "groups": [
            {
                "id": 1, "slug": "hunminjeongeum", "heritage": "훈민정음 해례본", "color": "#345D8C",
                "focus": "한글은 정말 세종 혼자 만들었을까?",
                "question": "훈민정음은 누가 만들고, 누가 설명을 붙였나요?",
                "source": "국가유산청 훈민정음 해례본·국사편찬위원회 자료",
                "items": [
                    ("세종은 1443년에 새 글자를 만들었습니다.", "○ 자료로 확인", "세종실록에 1443년 새 글자 창제 기록이 있습니다.", "훈민정음 1443 세종"),
                    ("1446년에는 새 글자의 원리와 사용법을 설명한 훈민정음 해례가 펴내졌습니다.", "○ 자료로 확인", "1446년 반포와 해례 간행이 확인됩니다.", "훈민정음 해례 1446"),
                    ("처음 만든 훈민정음 글자는 모두 28자였습니다.", "○ 자료로 확인", "처음 만든 글자는 초성 17자와 중성 11자입니다.", "훈민정음 처음 28자"),
                    ("집현전 학자들이 글자를 모두 만들었고 세종은 이름만 붙였습니다.", "× 자료와 다름", "글자 창제의 주체는 세종이며 학자들은 해설 등에 참여했습니다.", "세종 훈민정음 친제"),
                    ("훈민정음 해례본의 설명과 예시 작성에는 집현전 학자들이 참여했습니다.", "○ 자료로 확인", "정인지 등 집현전 학자들이 해례와 서문 작성에 참여했습니다.", "훈민정음 해례 집현전"),
                    ("각 글자를 정확히 어느 학자가 한 자씩 만들었는지 개인별 명단이 모두 남아 있습니다.", "× 자료와 다름", "글자별 담당자를 적은 개인 명단은 확인되지 않습니다.", "훈민정음 글자별 만든 사람"),
                ],
            },
            {
                "id": 2, "slug": "annals", "heritage": "조선왕조실록", "color": "#2F3436",
                "focus": "임금도 못 보는 기록이 어떻게 지켜졌나?",
                "question": "조선왕조실록은 누가 만들고, 어떻게 지켜졌나요?",
                "source": "국가유산청 조선왕조실록·국사편찬위원회 자료",
                "items": [
                    ("한 임금의 실록은 그 임금이 세상을 떠난 뒤 임시로 설치한 실록청에서 편찬했습니다.", "○ 자료로 확인", "실록 편찬 시기와 실록청 운영이 공식 자료에 나옵니다.", "실록 왕 사후 실록청"),
                    ("사관이 남긴 기록과 여러 자료를 모아 실록을 만들었습니다.", "○ 자료로 확인", "사초·시정기 등을 바탕으로 편찬했습니다.", "실록 사관 사초"),
                    ("조선의 임금은 매일 밤 자기 손으로 그날의 실록을 직접 썼습니다.", "× 자료와 다름", "사관의 기록을 바탕으로 후대에 편찬했습니다.", "조선왕조실록 누가 썼나"),
                    ("실록은 사고라는 여러 보관소에 나누어 두어 한 곳의 사고에 대비했습니다.", "○ 자료로 확인", "여러 사고에 복사본을 분산 보관했습니다.", "실록 사고 보관"),
                    ("임진왜란 때 전국의 실록이 모두 안전하게 남아 한 권도 잃지 않았습니다.", "× 자료와 다름", "여러 사고본이 소실되고 전주 사고본이 살아남았습니다.", "임진왜란 전주 사고본"),
                    ("전주 사고에 있던 실록은 전쟁 뒤 실록을 다시 만드는 중요한 바탕이 되었습니다.", "○ 자료로 확인", "전주 사고본을 토대로 여러 부를 다시 인쇄했습니다.", "전주 사고본 복원"),
                ],
            },
            {
                "id": 3, "slug": "hwaseong", "heritage": "수원 화성", "color": "#66737B",
                "focus": "거중기는 정말 공사를 바꿨을까?",
                "question": "수원 화성은 어떻게 지었고, 거중기는 어느 정도 역할을 했나요?",
                "source": "국가유산청 수원 화성·화성성역의궤 자료",
                "items": [
                    ("수원 화성은 정조 때인 1794년에 공사를 시작해 1796년에 완성되었습니다.", "○ 자료로 확인", "공사 기간은 공식 자료에서 확인됩니다.", "수원 화성 1794 1796"),
                    ("정약용은 화성 설계와 공사 방법을 마련하고 거중기 같은 도구 사용에 기여했습니다.", "○ 자료로 확인", "정약용의 설계와 기계 활용이 공식 안내에 설명되어 있습니다.", "정약용 화성 거중기"),
                    ("수원 화성은 거중기 한 대가 혼자 모든 돌을 옮겨 한 달 만에 완성했습니다.", "× 자료와 다름", "많은 인력·도구·공정이 필요했고 약 2년 반 동안 진행되었습니다.", "화성 공사 기간 인력"),
                    ("거중기 덕분에 공사 기간이 정확히 10년에서 34개월로 줄었다는 숫자는 당시 기록으로 완전히 증명되었습니다.", "△ 근거 부족·과장됨", "거중기의 기여는 확인되지만 특정 단축 수치의 출처를 따져야 합니다.", "거중기 기간 단축 근거"),
                    ("화성성역의궤에는 공사 과정, 재료, 인력과 비용에 관한 자세한 기록이 남아 있습니다.", "○ 자료로 확인", "의궤가 건축 과정을 상세히 전합니다.", "화성성역의궤 공사 기록"),
                    ("수원 화성은 조선과 관계없는 유럽 건축가가 혼자 설계했습니다.", "× 자료와 다름", "정조의 계획 아래 정약용 등 조선의 인물과 기술이 참여했습니다.", "수원 화성 설계"),
                ],
            },
            {
                "id": 4, "slug": "jagyeongnu-angbuilgu", "heritage": "자격루·앙부일구", "color": "#9C7432",
                "focus": "장영실은 그 뒤 어떻게 되었을까?",
                "question": "자격루와 앙부일구는 어떻게 시간을 알렸고, 장영실의 마지막 기록은 무엇인가요?",
                "source": "국사편찬위원회·국립고궁박물관 자료",
                "items": [
                    ("자격루는 물의 흐름과 기계 장치를 이용해 자동으로 시간을 알리는 물시계였습니다.", "○ 자료로 확인", "자동 물시계의 구조와 역할을 확인할 수 있습니다.", "자격루 자동 물시계"),
                    ("앙부일구는 해의 위치에 따라 생기는 그림자로 시간을 읽는 해시계입니다.", "○ 자료로 확인", "해 그림자와 시각선으로 시간을 읽습니다.", "앙부일구 해 그림자"),
                    ("장영실은 여러 과학 기구를 아무 도움 없이 혼자 설계하고 제작했습니다.", "× 자료와 다름", "세종의 지원 아래 여러 학자와 기술자가 함께 참여했습니다.", "장영실 공동 제작"),
                    ("앙부일구는 달빛만 있으면 밤에도 낮과 똑같이 시간을 정확히 읽을 수 있었습니다.", "× 자료와 다름", "해의 그림자를 이용하므로 햇빛이 필요합니다.", "앙부일구 밤 사용"),
                    ("장영실은 왕의 가마가 부서진 사건 뒤 벌을 받고 관직에서 물러났습니다.", "○ 자료로 확인", "세종실록에 가마 사건과 처벌 기록이 남아 있습니다.", "장영실 가마 사건"),
                    ("장영실은 그 뒤 제주도로 유배되어 1458년에 병으로 죽었다는 기록이 남아 있습니다.", "× 자료와 다름", "이후 행적과 사망 시기는 확실한 기록이 남아 있지 않습니다.", "장영실 이후 기록"),
                ],
            },
            {
                "id": 5, "slug": "jongmyo", "heritage": "종묘와 종묘제례악", "color": "#654C3E",
                "focus": "왜 이렇게 단순하게 지었을까?",
                "question": "종묘는 어떤 곳이며, 종묘제례악은 어떻게 이루어지나요?",
                "source": "국가유산청 종묘·국립무형유산원 자료",
                "items": [
                    ("종묘는 조선 왕과 왕비의 신주를 모시고 제사를 지내던 왕실 사당입니다.", "○ 자료로 확인", "종묘의 기능은 국가유산 자료에서 확인됩니다.", "종묘 왕실 사당"),
                    ("종묘제례악은 악기 연주뿐 아니라 노래와 춤이 함께하는 의례 음악입니다.", "○ 자료로 확인", "음악·노래·일무가 함께 진행됩니다.", "종묘제례악 노래 춤"),
                    ("종묘는 왕과 왕비가 평소 잠을 자며 생활하던 궁궐의 침전이었습니다.", "× 자료와 다름", "생활 공간이 아니라 조상 신주를 모시는 제례 공간입니다.", "종묘 사당 생활 공간"),
                    ("종묘제례악은 궁중 잔치의 흥을 돋우기 위해 매일 연주하던 오락 음악이었습니다.", "× 자료와 다름", "왕실 조상에게 제사를 지내는 엄숙한 의례 음악입니다.", "종묘제례악 언제 연주"),
                    ("종묘가 단순해 보이는 까닭은 오직 돈을 아끼기 위해서였다고 설계자의 기록에 적혀 있습니다.", "△ 근거 부족·과장됨", "제례의 엄숙함과 질서를 드러내며 비용 하나로 단정하기 어렵습니다.", "종묘 건축 단순한 이유"),
                    ("종묘제례악에는 북 소리만 있고 노래와 춤은 전혀 없습니다.", "× 자료와 다름", "악기·노래·춤이 함께하는 종합 의례입니다.", "종묘제례악 구성"),
                ],
            },
            {
                "id": 6, "slug": "nanjung-ilgi", "heritage": "난중일기", "color": "#25766C",
                "focus": "드라마 속 이순신은 진짜일까?",
                "question": "난중일기에는 이순신의 어떤 모습이 기록되어 있나요?",
                "source": "국가유산청 난중일기·국사편찬위원회 자료",
                "items": [
                    ("난중일기는 이순신이 임진왜란이 일어난 1592년부터 1598년까지 직접 쓴 일기입니다.", "○ 자료로 확인", "작성자와 작성 기간은 공식 자료에서 확인됩니다.", "난중일기 1592 1598"),
                    ("난중일기에는 전투뿐 아니라 날씨, 가족 걱정, 몸 상태와 개인적인 생각도 적혀 있습니다.", "○ 자료로 확인", "일기 원문에 일상과 감정 기록도 많습니다.", "난중일기 날씨 가족"),
                    ("난중일기는 전쟁이 모두 끝난 뒤 이순신이 기억을 떠올려 한 번에 쓴 책입니다.", "× 자료와 다름", "전쟁 중 날마다 또는 시기별로 남긴 일기입니다.", "난중일기 전쟁 중 작성"),
                    ("난중일기는 임금의 명령으로 작성한 공식 전투 보고서이므로 이순신의 개인 감정은 전혀 없습니다.", "× 자료와 다름", "개인 일기라서 감정과 일상이 함께 담겨 있습니다.", "난중일기 개인 감정"),
                    ("드라마에 나오는 이순신의 모든 대사는 난중일기 원문과 한 글자도 다르지 않습니다.", "△ 근거 부족·과장됨", "드라마에는 연출과 각색이 있으므로 장면마다 원문과 비교해야 합니다.", "드라마 난중일기 원문"),
                    ("난중일기에는 이순신이 노량해전에서 죽은 뒤 자신의 죽음을 직접 적은 마지막 장이 있습니다.", "× 자료와 다름", "자신의 죽음 뒤 기록을 직접 쓸 수 없습니다.", "난중일기 마지막 기록"),
                ],
            },
        ],
    },
}


def args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-root", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--font-dir", type=Path, default=DEFAULT_FONTS)
    # The 2025 classroom revision applies to three-kingdoms only; joseon materials are on
    # hold until that course is revised, so it is never re-rendered unless asked for.
    parser.add_argument("--eras", nargs="+", choices=sorted(DATA), default=["three-kingdoms"])
    return parser.parse_args()


def family(path: Path) -> str:
    font = FontToolsTTFont(str(path), lazy=True)
    try:
        for name_id in (16, 1):
            for record in font["name"].names:
                if record.nameID == name_id:
                    value = record.toUnicode().strip()
                    if value:
                        return value
    finally:
        font.close()
    return path.stem


def setup_fonts(font_dir: Path) -> tuple[str, str]:
    regular = font_dir / "SCDream5.ttf"
    heavy = font_dir / "SCDream9.ttf"
    if not regular.exists() or not heavy.exists():
        regular = Path("/usr/share/fonts/truetype/nanum/NanumGothic.ttf")
        heavy = Path("/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf")
    if not regular.exists() or not heavy.exists():
        raise FileNotFoundError("Korean font files are missing")
    pdfmetrics.registerFont(TTFont("SchoolRegular", str(regular)))
    pdfmetrics.registerFont(TTFont("SchoolHeavy", str(heavy)))
    return family(regular), family(heavy)


def p(text: str, style: ParagraphStyle) -> Paragraph:
    return Paragraph(str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), style)


def styles() -> dict[str, ParagraphStyle]:
    return {
        "body": ParagraphStyle("body", fontName="SchoolRegular", fontSize=7.1, leading=9.1, textColor=TEXT, wordWrap="CJK"),
        "small": ParagraphStyle("small", fontName="SchoolRegular", fontSize=6.4, leading=8.1, textColor=TEXT, wordWrap="CJK"),
        "tiny": ParagraphStyle("tiny", fontName="SchoolRegular", fontSize=5.7, leading=7.2, textColor=MUTED, wordWrap="CJK"),
        "white": ParagraphStyle("white", fontName="SchoolHeavy", fontSize=7.0, leading=8.4, textColor=WHITE, alignment=TA_CENTER, wordWrap="CJK"),
        "center": ParagraphStyle("center", fontName="SchoolRegular", fontSize=6.7, leading=8.2, textColor=TEXT, alignment=TA_CENTER, wordWrap="CJK"),
        "center_bold": ParagraphStyle("center_bold", fontName="SchoolHeavy", fontSize=6.8, leading=8.2, textColor=TEXT, alignment=TA_CENTER, wordWrap="CJK"),
        "answer": ParagraphStyle("answer", fontName="SchoolRegular", fontSize=6.1, leading=7.8, textColor=TEXT, wordWrap="CJK"),
    }


def draw_para(c: canvas.Canvas, text: str, style: ParagraphStyle, x: float, y_top: float, w: float, h: float) -> None:
    para = p(text, style)
    _, ph = para.wrap(w, h)
    para.drawOn(c, x, y_top - ph)


def header(c: canvas.Canvas, era: dict[str, Any], audience: str) -> float:
    width, height = A4
    x = 13 * mm
    c.setFillColor(PRIMARY)
    c.rect(0, height - 5 * mm, width, 5 * mm, stroke=0, fill=1)
    y = height - 18 * mm
    c.setFillColor(PRIMARY)
    c.roundRect(x, y, 21 * mm, 9 * mm, 3 * mm, stroke=0, fill=1)
    c.setFillColor(WHITE)
    c.setFont("SchoolHeavy", 9.2)
    c.drawCentredString(x + 10.5 * mm, y + 3 * mm, "2차시")
    c.setFillColor(TEXT)
    c.setFont("SchoolHeavy", 17.5)
    c.drawString(x + 27 * mm, y + 1.5 * mm, "AI에게 물어보았습니다")
    c.setFont("SchoolRegular", 7.0)
    c.setFillColor(MUTED)
    c.drawRightString(width - x, y + 3 * mm, f"AI혁신학교 공개수업 · {audience}")
    c.setStrokeColor(PRIMARY)
    c.setLineWidth(1.2)
    c.line(x, y - 2 * mm, width - x, y - 2 * mm)
    goal_y = y - 14 * mm
    c.setFillColor(MINT)
    c.rect(x, goal_y, width - 2 * x, 10 * mm, stroke=0, fill=1)
    c.setFillColor(PRIMARY)
    c.rect(x, goal_y, 2.2 * mm, 10 * mm, stroke=0, fill=1)
    c.setFont("SchoolHeavy", 8)
    c.drawString(x + 5 * mm, goal_y + 3.1 * mm, "학습목표")
    c.setFont("SchoolRegular", 7.4)
    c.setFillColor(TEXT)
    c.drawString(x + 26 * mm, goal_y + 3.1 * mm, "AI 설명에서 사실과 다른 부분이나 더 확인해야 할 부분을 찾고 이유를 설명할 수 있다.")
    return goal_y - 3.5 * mm


def group_box(c: canvas.Canvas, era: dict[str, Any], group: dict[str, Any], y_top: float, st: dict[str, ParagraphStyle]) -> float:
    width, _ = A4
    x = 13 * mm
    w = width - 26 * mm
    h = 27 * mm
    y = y_top - h
    c.setFillColor(WHITE)
    c.setStrokeColor(GRID)
    c.roundRect(x, y, w, h, 2.5 * mm, stroke=1, fill=1)
    accent = HexColor(group["color"])
    c.setFillColor(accent)
    c.roundRect(x, y, 42 * mm, h, 2.5 * mm, stroke=0, fill=1)
    c.setFillColor(WHITE)
    c.setFont("SchoolHeavy", 9)
    c.drawString(x + 4 * mm, y + h - 7 * mm, f'{group["id"]}모둠')
    draw_para(c, group["heritage"], ParagraphStyle("heritage", fontName="SchoolHeavy", fontSize=12, leading=14, textColor=WHITE, wordWrap="CJK"), x + 4 * mm, y + h - 10 * mm, 34 * mm, 14 * mm)
    c.setFont("SchoolRegular", 5.8)
    c.drawString(x + 4 * mm, y + 3 * mm, era["course"])
    rx = x + 46 * mm
    available = w - 46 * mm
    labels = [("핵심 질문", group["focus"]), ("AI 질문", group["question"]), ("학생", "5학년 ____반  이름 ____________________   역할  □자료  □기록  □기기  □발표")]
    row_h = 8.4 * mm
    for i, (label, value) in enumerate(labels):
        ry = y + h - (i + 1) * row_h
        c.setFillColor(MINT_LIGHT if i < 2 else WHITE)
        c.rect(rx, ry, available, row_h, stroke=0, fill=1)
        if i:
            c.setStrokeColor(GRID)
            c.line(rx, ry + row_h, x + w - 2 * mm, ry + row_h)
        c.setFont("SchoolHeavy", 6.8)
        c.setFillColor(PRIMARY_DARK)
        c.drawString(rx + 3 * mm, ry + 2.6 * mm, label)
        draw_para(c, value, st["small"], rx + 20 * mm, ry + row_h - 2.1 * mm, available - 23 * mm, row_h - 1 * mm)
    return y - 4 * mm


def section(c: canvas.Canvas, number: str, title: str, y_top: float) -> float:
    x = 13 * mm
    c.setFillColor(PRIMARY)
    c.roundRect(x, y_top - 7 * mm, 8 * mm, 7 * mm, 2 * mm, stroke=0, fill=1)
    c.setFillColor(WHITE)
    c.setFont("SchoolHeavy", 8.2)
    c.drawCentredString(x + 4 * mm, y_top - 4.8 * mm, number)
    c.setFillColor(TEXT)
    c.setFont("SchoolHeavy", 10.2)
    c.drawString(x + 11 * mm, y_top - 5.2 * mm, title)
    return y_top - 9.5 * mm


def legend_row(era: dict[str, Any]) -> list[str]:
    """The four judgement symbols. PPT, web app and worksheet must show this exact wording."""
    return ["○ 자료로 확인", "× 자료와 다름", f'△ {era["delta"]}', "? 더 찾아봐야 함"]


def form_styles() -> dict[str, ParagraphStyle]:
    """학교 제출 양식용. 작은 설명 글씨 없이 크게 읽히는 크기만 사용합니다."""
    return {
        "cell": ParagraphStyle("cell", fontName="SchoolRegular", fontSize=11, leading=15.5, textColor=FORM_INK, wordWrap="CJK"),
        "head": ParagraphStyle("head", fontName="SchoolHeavy", fontSize=11, leading=14, textColor=WHITE, alignment=TA_CENTER, wordWrap="CJK"),
        "mark": ParagraphStyle("mark", fontName="SchoolHeavy", fontSize=14, leading=18, textColor=FORM_INK, alignment=TA_CENTER, wordWrap="CJK"),
        "num": ParagraphStyle("num", fontName="SchoolHeavy", fontSize=12, leading=15, textColor=FORM_INK, alignment=TA_CENTER, wordWrap="CJK"),
        "legend": ParagraphStyle("legend", fontName="SchoolRegular", fontSize=10.5, leading=14, textColor=FORM_INK, alignment=TA_CENTER, wordWrap="CJK"),
        "label": ParagraphStyle("label", fontName="SchoolHeavy", fontSize=10.5, leading=14, textColor=FORM_INK, wordWrap="CJK"),
    }


def form_header(c: canvas.Canvas, era: dict[str, Any], lesson_no: str, title: str, objective: str) -> float:
    """‘삼국시대 N차시 + 이름란 + 학습 주제·목표’ 머리글. 1차시 활동지와 같은 모양입니다."""
    width, height = A4
    x = 14 * mm
    w = width - 28 * mm
    y = height - 26 * mm
    c.setFillColor(FORM_HEAD)
    c.setFont("SchoolHeavy", 24)
    c.drawString(x, y, f'{era["short"]} {lesson_no}')
    c.setFillColor(FORM_INK)
    c.setFont("SchoolRegular", 11)
    c.drawRightString(x + w, y + 2 * mm, "5학년 ______반   이름 ____________________")
    c.setStrokeColor(FORM_HEAD)
    c.setLineWidth(1.1)
    c.line(x, y - 4 * mm, x + w, y - 4 * mm)

    band_y = y - 26 * mm
    c.setFillColor(FORM_BAND)
    c.rect(x, band_y, w, 20 * mm, stroke=0, fill=1)
    c.setStrokeColor(FORM_LINE)
    c.setLineWidth(0.6)
    c.rect(x, band_y, w, 20 * mm, stroke=1, fill=0)
    c.line(x, band_y + 10 * mm, x + w, band_y + 10 * mm)
    c.line(x + 26 * mm, band_y, x + 26 * mm, band_y + 20 * mm)
    c.setFillColor(FORM_INK)
    c.setFont("SchoolHeavy", 10.5)
    c.drawString(x + 5 * mm, band_y + 13.5 * mm, "학습 주제")
    c.drawString(x + 5 * mm, band_y + 3.5 * mm, "학습 목표")
    c.setFont("SchoolHeavy", 13)
    c.drawString(x + 32 * mm, band_y + 13.2 * mm, title)
    c.setFont("SchoolRegular", 10.5)
    c.drawString(x + 32 * mm, band_y + 3.4 * mm, objective)
    return band_y - 7 * mm


def form_footer(c: canvas.Canvas) -> None:
    width, _ = A4
    x = 14 * mm
    c.setFillColor(FORM_SOFT)
    c.setFont("SchoolRegular", 8)
    c.drawString(x, 8 * mm, "인공지능과 역사")
    c.drawRightString(width - x, 8 * mm, "S-Core Dream")


def student_page(path: Path, era: dict[str, Any], group: dict[str, Any]) -> None:
    st = form_styles()
    path.parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(str(path), pagesize=A4)
    c.setTitle(f'{era["short"]} 2차시 {group["id"]}모둠 {group["heritage"]} 활동지')
    c.setAuthor("MOAKIT")
    width, _ = A4
    x = 14 * mm
    w = width - 28 * mm
    y = form_header(c, era, "2차시", "AI에게 물어보았습니다", "AI 설명을 ○×△?로 판단하고 공식 자료의 출처로 확인할 수 있다.")

    group_row = Table(
        [[p("모둠·유산", st["label"]), p(f'{group["id"]}모둠 · {group["heritage"]}', st["cell"]),
          p("역할", st["label"]), p("□ 자료   □ 기록   □ 기기   □ 발표", st["cell"])]],
        colWidths=[22 * mm, 62 * mm, 16 * mm, w - 100 * mm],
        rowHeights=[13 * mm],
    )
    group_row.setStyle(TableStyle([
        ("GRID", (0, 0), (-1, -1), 0.6, FORM_LINE), ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6), ("RIGHTPADDING", (0, 0), (-1, -1), 6),
    ]))
    group_row.wrapOn(c, w, 13 * mm); group_row.drawOn(c, x, y - 13 * mm)
    y -= 19 * mm

    rows = [[p(h, st["head"]) for h in ["번호", "AI가 한 말", "내 판단 (○×△?)", "확인한 출처"]]]
    for i, item in enumerate(group["items"], 1):
        rows.append([p(str(i), st["num"]), p(item[0], st["cell"]), p("○   ×   △   ?", st["mark"]), p("", st["cell"])])
    heights = [11 * mm] + [24 * mm] * 6
    table = Table(rows, colWidths=[12 * mm, w - 88 * mm, 36 * mm, 40 * mm], rowHeights=heights)
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), FORM_HEAD), ("GRID", (0, 0), (-1, -1), 0.5, FORM_LINE),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
    ]))
    th = sum(heights)
    table.wrapOn(c, w, th); table.drawOn(c, x, y - th)
    y -= th + 5 * mm

    legend = Table([[p(label, st["legend"]) for label in legend_row(era)]], colWidths=[w / 4] * 4, rowHeights=[12 * mm])
    legend.setStyle(TableStyle([
        ("GRID", (0, 0), (-1, -1), 0.5, FORM_LINE), ("BACKGROUND", (0, 0), (-1, -1), FORM_BAND),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    legend.wrapOn(c, w, 12 * mm); legend.drawOn(c, x, y - 12 * mm)
    y -= 17 * mm

    c.setStrokeColor(FORM_LINE); c.setLineWidth(0.6)
    c.rect(x, y - 18 * mm, w, 18 * mm, stroke=1, fill=0)
    c.setFillColor(FORM_INK)
    c.setFont("SchoolHeavy", 11)
    c.drawString(x + 6 * mm, y - 11 * mm, "오늘의 한 문장")
    c.setFont("SchoolRegular", 11)
    c.drawString(x + 40 * mm, y - 11 * mm, "AI의 말은 ____________________________________ 때문에 확인해야 한다.")

    form_footer(c)
    c.showPage(); c.save()


def merge(paths: Iterable[Path], output: Path) -> None:
    writer = PdfWriter()
    for path in paths:
        for page in PdfReader(str(path)).pages:
            writer.add_page(page)
    with output.open("wb") as f:
        writer.write(f)


def teacher_pdf(path: Path, era: dict[str, Any], answers_only: bool = False) -> None:
    st = styles(); c = canvas.Canvas(str(path), pagesize=A4); c.setAuthor("MOAKIT")
    width, _ = A4; x = 13 * mm; w = width - 26 * mm
    if not answers_only:
        y = header(c, era, "교사용 운영안")
        c.setFillColor(TEXT); c.setFont("SchoolHeavy", 17); c.drawString(x, y - 7 * mm, f'{era["short"]} 2차시 운영안')
        c.setFillColor(MUTED); c.setFont("SchoolRegular", 7.5); c.drawString(x, y - 14 * mm, "PPT·활동 화면·활동지가 같은 항목명(내 판단 ○×△? · 확인한 출처 · 오늘의 한 문장)을 사용합니다.")
        rows = [[p(h, st["white"]) for h in ["시간", "교수·학습 활동", "교사 포인트"]]]
        schedule = [
            ("5분", "AI 신뢰도 손들기(PPT 2)와 오늘의 미션 안내(PPT 3~4)", "정답 슬라이드를 먼저 열지 않습니다."),
            ("8분", "활동지 ‘내 판단’ 칸에 6문장을 ○×△?로 표시", "맞는 말과 틀린 말이 섞여 있음을 강조합니다."),
            ("12분", "모둠별로 공식 자료를 확인하고 ‘확인한 출처’ 칸을 채움", "핵심 낱말 2~4개로 검색하게 합니다(PPT 7)."),
            ("7분", "모둠 발표: 문장 번호 → 내 판단 → 확인한 출처", "판단과 함께 확인한 기관 이름을 말하게 합니다."),
            ("8분", "발표한 모둠부터 정답 슬라이드 공개, ‘오늘의 한 문장’ 쓰기", "PPT 8~13은 해당 모둠 발표 뒤에만 엽니다."),
        ]
        for row in schedule:
            rows.append([p(row[0], st["center_bold"]), p(row[1], st["body"]), p(row[2], st["small"])])
        heights = [9 * mm] + [20 * mm] * 5
        t = Table(rows, colWidths=[18 * mm, 98 * mm, w - 116 * mm], rowHeights=heights)
        t.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), PRIMARY_DARK), ("GRID", (0, 0), (-1, -1), 0.5, GRID), ("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("LEFTPADDING", (0, 0), (-1, -1), 5), ("RIGHTPADDING", (0, 0), (-1, -1), 5), ("BACKGROUND", (0, 2), (-1, 2), MINT_LIGHT), ("BACKGROUND", (0, 4), (-1, 4), MINT_LIGHT)]))
        th = sum(heights); t.wrapOn(c, w, th); t.drawOn(c, x, y - 21 * mm - th)
        ny = y - 29 * mm - th
        c.setFillColor(BEIGE); c.roundRect(x, ny - 40 * mm, w, 40 * mm, 2 * mm, stroke=0, fill=1)
        c.setFillColor(PRIMARY_DARK); c.setFont("SchoolHeavy", 8); c.drawString(x + 4 * mm, ny - 7 * mm, "학교 양식 적용 기준")
        bullets = ["학생 활동지는 모둠별 A4 한 장, 한 페이지로만 인쇄합니다.", "활동지는 6문장 판단과 확인한 출처, 오늘의 한 문장까지만 기록합니다.", "‘?’는 아직 확인하지 못했다는 뜻이며 틀린 답이 아닙니다. 정답표에는 나오지 않습니다.", "완성한 기록지는 10차시 교실 박물관의 검증 증거물로 전시합니다."]
        for i, text in enumerate(bullets):
            c.setFillColor(TEXT); c.setFont("SchoolRegular", 7); c.drawString(x + 5 * mm, ny - (14 + i * 6) * mm, f"• {text}")
        c.setFillColor(MUTED); c.setFont("SchoolRegular", 5.5); c.drawRightString(width - x, 7 * mm, f'{era["short"]} · 교사용 · S-Core Dream')
        c.showPage()
    for group in era["groups"]:
        y = header(c, era, "교사용 정답")
        c.setFillColor(HexColor(group["color"])); c.roundRect(x, y - 15 * mm, 28 * mm, 11 * mm, 3 * mm, stroke=0, fill=1)
        c.setFillColor(WHITE); c.setFont("SchoolHeavy", 9); c.drawCentredString(x + 14 * mm, y - 11.5 * mm, f'{group["id"]}모둠')
        c.setFillColor(TEXT); c.setFont("SchoolHeavy", 15); c.drawString(x + 34 * mm, y - 12 * mm, group["heritage"])
        c.setFont("SchoolRegular", 7); c.setFillColor(MUTED); c.drawRightString(width - x, y - 11.5 * mm, group["focus"])
        rows = [[p(h, st["white"]) for h in ["번호", "AI가 한 말", "판정", "교사용 근거"]]]
        for i, item in enumerate(group["items"], 1):
            rows.append([p(str(i), st["center_bold"]), p(item[0], st["answer"]), p(item[1], st["center_bold"]), p(item[2], st["answer"])])
        heights = [9 * mm] + [24 * mm] * 6
        t = Table(rows, colWidths=[10 * mm, 80 * mm, 31 * mm, w - 121 * mm], rowHeights=heights)
        t.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), PRIMARY_DARK), ("GRID", (0, 0), (-1, -1), 0.45, GRID), ("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("LEFTPADDING", (0, 0), (-1, -1), 4), ("RIGHTPADDING", (0, 0), (-1, -1), 4), ("BACKGROUND", (0, 2), (-1, 2), MINT_LIGHT), ("BACKGROUND", (0, 4), (-1, 4), MINT_LIGHT), ("BACKGROUND", (0, 6), (-1, 6), MINT_LIGHT)]))
        th = sum(heights); t.wrapOn(c, w, th); t.drawOn(c, x, y - 20 * mm - th)
        sy = y - 25 * mm - th
        c.setFillColor(BEIGE); c.roundRect(x, sy - 34 * mm, w, 34 * mm, 2 * mm, stroke=0, fill=1)
        c.setFillColor(PRIMARY_DARK); c.setFont("SchoolHeavy", 7.3); c.drawString(x + 4 * mm, sy - 6 * mm, "추천 검색어")
        for i, item in enumerate(group["items"], 1):
            col = 0 if i <= 3 else 1; row = i - 1 if i <= 3 else i - 4
            c.setFillColor(TEXT); c.setFont("SchoolRegular", 6.2); c.drawString(x + 4 * mm + col * w / 2, sy - (13 + row * 7) * mm, f"{i}. {item[3]}")
        c.setFillColor(MUTED); c.setFont("SchoolRegular", 5.5); c.drawString(x + 4 * mm, sy - 31 * mm, f'우선 확인처: {group["source"]}')
        c.drawRightString(width - x, 7 * mm, f'{era["short"]} · {group["id"]}모둠 · 정답 · S-Core Dream')
        c.showPage()
    c.save()


def add_text(slide, x: float, y: float, w: float, h: float, text: str, font: str, size: float, color: RGBColor = PPT_TEXT, bold: bool = False, align=PP_ALIGN.LEFT, fill: RGBColor | None = None) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame; tf.clear(); tf.margin_left = Inches(0.07); tf.margin_right = Inches(0.07); tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]; para.alignment = align; para.text = text
    for run in para.runs:
        run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color


def rect(slide, x: float, y: float, w: float, h: float, fill: RGBColor, rounded: bool = False) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()


def ppt_header(slide, era: dict[str, Any], regular: str, heavy: str) -> None:
    rect(slide, 0, 0, 13.333, 0.12, PPT_PRIMARY)
    rect(slide, 0.55, 0.36, 1.15, 0.42, PPT_PRIMARY, True)
    add_text(slide, 0.55, 0.36, 1.15, 0.42, "2차시", heavy, 12, RGBColor(255,255,255), True, PP_ALIGN.CENTER)
    add_text(slide, 1.88, 0.28, 5.8, 0.58, "AI에게 물어보았습니다", heavy, 25, PPT_TEXT, True)
    add_text(slide, 9.4, 0.35, 3.35, 0.35, f'{era["short"]} · 교사용 수업 PPT', regular, 9, PPT_MUTED, False, PP_ALIGN.RIGHT)
    rect(slide, 0.55, 0.93, 12.23, 0.025, PPT_PRIMARY)


def ppt(path: Path, era: dict[str, Any], regular: str, heavy: str) -> None:
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5); blank = prs.slide_layouts[6]
    def new(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(blank); ppt_header(slide, era, regular, heavy)
        add_text(slide, 0.72, 1.18, 11.9, 0.62, title, heavy, 28, PPT_TEXT, True)
        if subtitle: add_text(slide, 0.75, 1.77, 11.8, 0.38, subtitle, regular, 12, PPT_MUTED)
        return slide
    white = RGBColor(255, 255, 255)

    # 1 · 표지
    slide = prs.slides.add_slide(blank); rect(slide, 0, 0, 13.333, 7.5, RGBColor(248,251,249)); rect(slide, 0, 0, 13.333, 0.2, PPT_PRIMARY)
    add_text(slide, 0.82, 0.75, 6, 0.36, "AI혁신학교 공개수업 · 2차시", regular, 12, PPT_DARK, True)
    add_text(slide, 0.82, 1.38, 10.8, 1.18, "AI에게 물어보았습니다", heavy, 42, PPT_TEXT, True)
    rect(slide, 0.82, 3.05, 11.55, 1.38, PPT_MINT); rect(slide, 0.82, 3.05, 0.1, 1.38, PPT_PRIMARY)
    add_text(slide, 1.18, 3.28, 10.7, 0.35, "핵심 질문", heavy, 13, PPT_DARK, True)
    add_text(slide, 1.18, 3.68, 10.7, 0.55, f'“{era["core"]}”', heavy, 22, PPT_DARK, True)
    add_text(slide, 0.82, 5.05, 11.5, 0.95, "우리 모둠 유산에 대한 AI 문장 6개에서\n의심할 문장을 찾아 출처로 확인합니다.", regular, 21)

    # 2 · 도입
    slide = new("AI의 역사 설명을 얼마나 믿나요?", "손을 들어 봅시다. 정답은 없습니다.")
    for i, text in enumerate(["거의 다 맞을 것 같다", "반쯤은 맞을 것 같다", "그대로 믿기는 어렵다"]):
        x = 0.85 + i * 4.08
        rect(slide, x, 2.35, 3.75, 2.15, PPT_MINT if i != 2 else PPT_BEIGE, True)
        add_text(slide, x, 2.35, 3.75, 2.15, text, heavy, 21, PPT_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, 0.85, 5.15, 11.6, 1.0, "오늘은 AI가 우리 모둠 유산을 설명한 6문장을 직접 확인합니다.", heavy, 22, PPT_TEXT, True, PP_ALIGN.CENTER)

    # 3 · 오늘의 미션
    slide = new("오늘의 미션", "우리 모둠 유산 6문장에서 의심 문장 찾기")
    rect(slide, 0.85, 2.25, 11.6, 1.5, PPT_PRIMARY, True)
    add_text(slide, 1.1, 2.25, 11.1, 1.5, "우리 모둠 유산 6문장에서 의심 문장을 찾아라", heavy, 30, white, True, PP_ALIGN.CENTER)
    for i, text in enumerate(["검색하기 전에 먼저 내 판단을 표시한다.", "맞는 말과 틀린 말이 섞여 있음을 기억한다.", "판단한 뒤에는 확인한 출처를 꼭 적는다."]):
        rect(slide, 1.0, 4.15 + i * 0.92, 0.58, 0.58, PPT_PRIMARY, True)
        add_text(slide, 1.0, 4.15 + i * 0.92, 0.58, 0.58, str(i+1), heavy, 17, white, True, PP_ALIGN.CENTER)
        add_text(slide, 1.85, 4.12 + i * 0.92, 10.3, 0.65, text, regular, 20)

    # 4 · 판단 기호 (활동지·활동 화면과 같은 표현)
    slide = new("내 판단 (○×△?)", "활동지의 ‘내 판단’ 칸에 그대로 표시합니다.")
    for i, label in enumerate(legend_row(era)):
        symbol, meaning = label.split(" ", 1)
        x = 0.85 + i * 3.05
        rect(slide, x, 2.18, 2.65, 2.55, PPT_MINT if i % 2 == 0 else RGBColor(247,249,248), True)
        add_text(slide, x, 2.35, 2.65, 1.0, symbol, heavy, 38, PPT_DARK, True, PP_ALIGN.CENTER)
        add_text(slide, x+0.12, 3.5, 2.4, 0.85, meaning, regular, 15, PPT_TEXT, False, PP_ALIGN.CENTER)
    rect(slide, 1.2, 5.25, 10.95, 0.75, PPT_BEIGE, True)
    add_text(slide, 1.35, 5.25, 10.65, 0.75, "△와 ?도 정답이 될 수 있습니다. 모르는 것을 억지로 채우지 않습니다.", heavy, 17, PPT_DARK, True, PP_ALIGN.CENTER)

    # 5 · 활동 화면 안내
    slide = new("활동 화면 여는 방법", "로그인 없이 우리 모둠만 고르면 바로 시작합니다.")
    rect(slide, 0.85, 2.2, 11.6, 1.35, PPT_MINT, True); rect(slide, 0.85, 2.2, 0.12, 1.35, PPT_PRIMARY)
    add_text(slide, 1.2, 2.2, 11.1, 1.35, f'수업 사이트 → {era["short"]} → 2차시 → ‘{era["appTab"]}’ 탭', heavy, 24, PPT_DARK, True)
    add_text(slide, 1.2, 3.62, 11.1, 0.45, f'주소 뒤에 {era["appPath"]} 를 붙여도 바로 열립니다.', regular, 14, PPT_MUTED)
    for i, text in enumerate(["학생 로그인이 없습니다. 우리 모둠만 고르면 됩니다.", "이 화면은 내 판단만 저장하고 정답·점수는 보여 주지 않습니다.", "종이 활동지와 화면의 문장 번호는 똑같습니다."]):
        rect(slide, 1.0, 4.35 + i * 0.82, 0.12, 0.12, PPT_PRIMARY, True)
        add_text(slide, 1.3, 4.25 + i * 0.82, 10.9, 0.65, text, regular, 19)

    # 6 · 활동 시간 안내
    slide = new("활동 시간 안내", "35분을 네 덩어리로 나눕니다.")
    steps = [
        ("8분", "혼자 판단", "6문장을 ○×△?로 표시합니다."),
        ("12분", "모둠 조사", "공식 자료로 확인하고 ‘확인한 출처’ 칸을 채웁니다."),
        ("7분", "모둠 발표", "문장 번호 → 내 판단 → 확인한 출처 순서로 말합니다."),
        ("8분", "정답 확인·정리", "발표한 모둠부터 정답을 보고 ‘오늘의 한 문장’을 씁니다."),
    ]
    for i, (minutes, title, detail) in enumerate(steps):
        y = 2.2 + i * 1.12
        rect(slide, 0.85, y, 1.5, 0.88, PPT_PRIMARY, True)
        add_text(slide, 0.85, y, 1.5, 0.88, minutes, heavy, 22, white, True, PP_ALIGN.CENTER)
        add_text(slide, 2.6, y - 0.04, 3.0, 0.5, title, heavy, 20, PPT_TEXT, True)
        add_text(slide, 2.6, y + 0.42, 9.6, 0.48, detail, regular, 15, PPT_MUTED)

    # 7 · 좋은 검색어 + 어디에서 확인할까
    slide = new("좋은 검색어 만들기 · 어디에서 확인할까?")
    add_text(slide, 0.85, 2.12, 5.7, 0.4, "좋은 검색어 만들기", heavy, 18, PPT_DARK, True)
    rect(slide, 0.85, 2.62, 5.7, 1.15, RGBColor(252,239,239), True)
    add_text(slide, 1.05, 2.62, 5.3, 1.15, "✕  AI 문장 전체를 그대로 붙여넣기", regular, 16, RGBColor(154,62,62))
    rect(slide, 0.85, 3.92, 5.7, 1.15, PPT_MINT, True)
    add_text(slide, 1.05, 3.92, 5.3, 1.15, "○  핵심 낱말 2~4개 + 확인할 쟁점", heavy, 16, PPT_DARK, True)
    add_text(slide, 0.85, 5.22, 5.7, 0.9, "문장 전체가 아니라 확인할 핵심만 남깁니다.", regular, 14, PPT_MUTED)
    add_text(slide, 6.95, 2.12, 5.5, 0.4, "어디에서 확인할까? (공식 자료 순위)", heavy, 18, PPT_DARK, True)
    for i, text in enumerate(era["sourceRanks"]):
        rect(slide, 6.95, 2.62 + i * 0.78, 5.5, 0.62, PPT_PRIMARY if i < 3 else PPT_BEIGE, True)
        add_text(slide, 7.15, 2.62 + i * 0.78, 5.1, 0.62, text, heavy if i < 3 else regular, 14, white if i < 3 else PPT_TEXT, i < 3)
    add_text(slide, 6.95, 5.72, 5.5, 0.5, "확인한 곳의 기관 이름을 활동지 ‘확인한 출처’ 칸에 적습니다.", regular, 13, PPT_MUTED)

    # 8~13 · 모둠별 정답 (발표 뒤에만 공개)
    for group in era["groups"]:
        slide = new(f'{group["id"]}모둠 정답 · {group["heritage"]}', "이 모둠의 발표를 들은 뒤에 공개하세요.")
        rect(slide, 10.75, 0.3, 1.95, 0.52, RGBColor(161,73,73), True)
        add_text(slide, 10.75, 0.3, 1.95, 0.52, "발표 후 공개", heavy, 12, white, True, PP_ALIGN.CENTER)
        add_text(slide, 1.82, 2.16, 4.3, 0.3, "AI가 한 말", heavy, 11, PPT_MUTED)
        add_text(slide, 6.32, 2.16, 3.85, 0.3, "왜 그렇게 판단하나요", heavy, 11, PPT_MUTED)
        add_text(slide, 10.35, 2.16, 2.55, 0.3, "추천 검색어", heavy, 11, PPT_MUTED)
        color = RGBColor(*tuple(int(group["color"][j:j+2], 16) for j in (1,3,5)))
        rect(slide, 0.82, 2.52, 12.05, 0.02, color)
        for i, item in enumerate(group["items"], 1):
            y = 2.62 + (i-1) * 0.74
            verdict_color = PPT_PRIMARY if item[1].startswith("○") else RGBColor(161,73,73) if item[1].startswith("×") else RGBColor(161,122,47)
            rect(slide, 0.82, y, 0.86, 0.5, verdict_color, True)
            add_text(slide, 0.82, y, 0.86, 0.5, f'{i} {item[1].split()[0]}', heavy, 12, white, True, PP_ALIGN.CENTER)
            add_text(slide, 1.82, y-0.03, 4.3, 0.58, item[0], regular, 10.4)
            add_text(slide, 6.32, y-0.03, 3.85, 0.58, item[2], regular, 10.2, PPT_MUTED)
            add_text(slide, 10.35, y-0.03, 2.55, 0.58, item[3], regular, 9.8, PPT_DARK)
        add_text(slide, 0.82, 7.06, 12.05, 0.34, f'우선 확인처: {group["source"]}', regular, 10, PPT_MUTED)

    # 14 · 정리 + 다음 차시 예고
    slide = new("AI는 왜 틀릴까?", "그리고 다음 시간에는 무엇을 할까요?")
    for i, text in enumerate(["AI는 여러 글의 표현과 패턴을 조합해 답합니다.", "자료에 빈칸이 있어도 그럴듯하게 이어 말합니다.", "자신 있는 말투가 사실을 보장하지 않습니다.", "출처·근거·원본은 사람이 다시 확인해야 합니다."]):
        rect(slide, 1.0, 2.28 + i * 0.72, 0.12, 0.12, PPT_PRIMARY, True)
        add_text(slide, 1.3, 2.18 + i * 0.72, 10.9, 0.6, text, regular, 19)
    add_text(slide, 0.85, 5.15, 11.6, 0.4, "다음 시간 · 오늘 우리가 한 확인에 이름 붙이기", heavy, 16, PPT_DARK, True)
    ranks = era["next"].split(" · ")
    for i, step in enumerate(ranks):
        width = 11.3 / len(ranks); x = 0.85 + i * (11.7 / len(ranks))
        special = i == len(ranks) - 1 and era["short"] == "삼국시대"
        rect(slide, x, 5.62, width, 1.0, PPT_PRIMARY if special else PPT_MINT, True)
        add_text(slide, x, 5.62, width, 1.0, step, heavy, 19, white if special else PPT_DARK, True, PP_ALIGN.CENTER)
    path.parent.mkdir(parents=True, exist_ok=True); prs.save(str(path))


def current_lesson_ids(output: Path, era_id: str) -> set[int]:
    """manifest에 남아 있는 차시만 시대 묶음에 넣는다. 과정에서 빠진 차시의 옛 파일은 제외."""
    path = output / "manifest.json"
    if not path.exists():
        return set()
    data = json.loads(path.read_text(encoding="utf-8"))
    return {item["lessonId"] for item in data.get("eras", {}).get(era_id, {}).get("lessons", [])}


def in_current_course(filename: str, lesson_ids: set[int]) -> bool:
    match = re.match(r"lesson-(\d{2})-", filename)
    return True if not match else int(match.group(1)) in lesson_ids


def rebuild(zip_path: Path, files: Iterable[Path], base: Path) -> None:
    temp = zip_path.with_suffix(".tmp.zip")
    with zipfile.ZipFile(temp, "w", zipfile.ZIP_DEFLATED, compresslevel=9) as archive:
        for file in sorted({f for f in files if f.exists()}):
            if file.resolve() in {zip_path.resolve(), temp.resolve()}:
                continue
            archive.write(file, file.relative_to(base))
    temp.replace(zip_path)


def update_manifest(output: Path) -> None:
    path = output / "manifest.json"
    if not path.exists():
        return
    data = json.loads(path.read_text(encoding="utf-8"))
    for era_id, era in DATA.items():
        era_record = data.get("eras", {}).get(era_id)
        if not era_record:
            continue
        lesson = next((item for item in era_record.get("lessons", []) if item.get("lessonId") == 2), None)
        if not lesson:
            continue
        files = lesson.setdefault("files", {})
        for key, name in [("student", "lesson-02-student.pdf"), ("teacher", "lesson-02-teacher.pdf"), ("answer", "lesson-02-answer.pdf"), ("bundle", "lesson-02-all.zip"), ("ppt", "lesson-02-teaching.pptx")]:
            file = output / era_id / name
            if file.exists():
                files[key] = {"path": f"/downloads/{era_id}/{name}", "size": file.stat().st_size}
        files["groups"] = [{"groupId": g["id"], "heritage": g["heritage"], "path": f'/downloads/{era_id}/lesson-02-group-{g["id"]:02d}-{g["slug"]}.pdf', "size": (output / era_id / f'lesson-02-group-{g["id"]:02d}-{g["slug"]}.pdf').stat().st_size, "pages": 1, "format": "A4"} for g in era["groups"]]
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def lesson_two_files(directory: Path, era: dict[str, Any]) -> list[Path]:
    names = ["lesson-02-student.pdf", "lesson-02-teacher.pdf", "lesson-02-answer.pdf", "lesson-02-teaching.pptx", "lesson-02-all.zip"]
    names += [f'lesson-02-group-{group["id"]:02d}-{group["slug"]}.pdf' for group in era["groups"]]
    return [directory / name for name in names]


def main() -> None:
    options = args(); output = options.output_root; regular_name, heavy_name = setup_fonts(options.font_dir); output.mkdir(parents=True, exist_ok=True)
    for era_id in options.eras:
        era = DATA[era_id]
        directory = output / era_id; directory.mkdir(parents=True, exist_ok=True); groups: list[Path] = []
        for group in era["groups"]:
            file = directory / f'lesson-02-group-{group["id"]:02d}-{group["slug"]}.pdf'; student_page(file, era, group); groups.append(file)
        student = directory / "lesson-02-student.pdf"; merge(groups, student)
        teacher = directory / "lesson-02-teacher.pdf"; teacher_pdf(teacher, era, False)
        answer = directory / "lesson-02-answer.pdf"; teacher_pdf(answer, era, True)
        deck = directory / "lesson-02-teaching.pptx"; ppt(deck, era, regular_name, heavy_name)
        lesson_zip = directory / "lesson-02-all.zip"; rebuild(lesson_zip, [student, teacher, answer, deck, *groups], directory)
        era_zip = directory / f"{era_id}-all-materials.zip"
        lesson_ids = current_lesson_ids(output, era_id)
        rebuild(era_zip, [f for f in directory.iterdir() if f.is_file() and f.name != era_zip.name and in_current_course(f.name, lesson_ids)], directory)
    update_manifest(output)
    # The cross-era bundle always ships every era's current files, including the ones this run left alone.
    complete = [file for era_id, era in DATA.items() for file in lesson_two_files(output / era_id, era)]
    rebuild(output / "lesson-02-samguk-joseon-complete.zip", complete, output)
    (output / "lesson-02-materials.json").write_text(json.dumps({"format": "A4 portrait", "student_pages_per_group": 1, "font": "S-Core Dream" if "S-Core" in regular_name else regular_name, "eras": {era_id: {"groups": 6, "student_pages": 6, "ppt": f"/{era_id}/lesson-02-teaching.pptx"} for era_id in DATA}}, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(f"Generated one-page lesson 2 materials in {output}")


if __name__ == "__main__":
    main()
