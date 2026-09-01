"""4차시 ‘데이터 만들기’ 인쇄 자료 생성기.

여섯 모둠이 똑같은 일곱 항목을 채워야 학급 데이터 표가 만들어집니다.
PPT · 웹앱(src/content/three-kingdoms/webActivities.ts) · 활동지가 아래 FIELDS의
항목 이름을 토씨 하나 다르지 않게 함께 사용합니다.

2차시 자료와 같은 학교 양식(초록·민트·베이지 + S-Core Dream)을 씁니다.
"""

from __future__ import annotations

import argparse
import json
import zipfile
from pathlib import Path
from typing import Any, Iterable

from fontTools.ttLib import TTFont as FontToolsTTFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pypdf import PdfReader, PdfWriter
from reportlab.lib import colors
from reportlab.lib.colors import HexColor
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph, Table, TableStyle

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = ROOT / "public" / "downloads"
DEFAULT_FONTS = ROOT / "public" / "fonts"

PRIMARY = HexColor("#2F715B")
PRIMARY_DARK = HexColor("#1F4F40")
MINT = HexColor("#EAF3EF")
MINT_LIGHT = HexColor("#F6FAF8")
BEIGE = HexColor("#FBF2E2")
GRID = HexColor("#C8D2CE")
TEXT = HexColor("#1F2925")
MUTED = HexColor("#64716C")
WHITE = colors.white

PPT_PRIMARY = RGBColor(47, 113, 91)
PPT_DARK = RGBColor(31, 79, 64)
PPT_MINT = RGBColor(234, 243, 239)
PPT_BEIGE = RGBColor(251, 242, 226)
PPT_TEXT = RGBColor(31, 41, 37)
PPT_MUTED = RGBColor(100, 113, 108)
PPT_WHITE = RGBColor(255, 255, 255)

# 활동지·PPT·웹앱이 함께 쓰는 일곱 항목. 순서까지 같아야 학급 표의 열이 맞습니다.
FIELDS: list[tuple[str, str]] = [
    ("시기", "언제 만들었는지 (세기·왕 이름)"),
    ("만든 까닭", "무엇을 위해 만들었는지"),
    ("가치", "왜 중요한 유산인지"),
    ("현재 상태", "지금 어디에 어떤 모습으로 남아 있는지"),
    ("AI 오류 바로잡기", "2차시에서 ×였던 문장을 바르게 고쳐 쓰기"),
    ("아직 모름", "자료로 확인되지 않은 점"),
    ("출처", "확인한 기관 이름"),
]

SOURCE_RANKS = [
    "1순위  국가유산청 국가유산포털",
    "2순위  국립중앙박물관·국립부여박물관",
    "3순위  유네스코 세계유산센터·우리역사넷",
    "비교용  블로그·영상·AI 답변",
]

# 2차시에 학생들이 붙인 확인 방법의 이름. 4차시 도입 5분에 그대로 씁니다.
VERIFICATION_NAMES = [
    ("출처", "누가 만든 자료인지 확인했다"),
    ("시기", "언제 만든 자료인지 확인했다"),
    ("교차", "다른 자료와 견주어 보았다"),
    ("원본", "요약문 대신 원본까지 갔다"),
    ("보류", "확인되지 않으면 그대로 두었다"),
]

ERA = {
    "id": "three-kingdoms",
    "short": "삼국시대",
    "course": "삼국시대 문화유산 해설사",
    "core": "여섯 유산을 어떻게 하나의 표로 비교할까?",
    "appTab": "활동지 수업",
    "appPath": "/three-kingdoms/lesson/4?view=activity",
}

# 모범 답안은 교사용 PDF·PPT에만 실립니다. 학생 화면과 활동지에는 넣지 않습니다.
GROUPS: list[dict[str, Any]] = [
    {
        "id": 1, "slug": "muryeongwangneung", "heritage": "무령왕릉", "color": "#7A5AA6",
        "focus": "발굴은 왜 아쉬움으로 남았을까?",
        "source": "국립공주박물관·국립중앙박물관 무령왕릉 자료",
        "answers": [
            "백제 무령왕 때, 6세기 초",
            "무령왕과 왕비를 모시기 위해 만든 벽돌무덤",
            "지석 덕분에 무덤 주인과 시기를 아는 삼국시대 왕릉",
            "공주 무령왕릉과 왕릉원에 남아 있고 유물은 국립공주박물관 소장",
            "‘도굴되어 유물이 사라졌다’ → 도굴되지 않아 많은 유물이 함께 나왔다",
            "유물이 놓인 까닭과 장례의 모든 과정",
            "국립공주박물관·국립중앙박물관",
        ],
    },
    {
        "id": 2, "slug": "baekje-incense-burner", "heritage": "백제 금동대향로", "color": "#3D7F5D",
        "focus": "향로에 새겨진 것은 무엇을 뜻할까?",
        "source": "국립부여박물관 백제 금동대향로 자료",
        "answers": [
            "백제, 6~7세기 무렵 (1993년 출토)",
            "절에서 의례에 쓰기 위해 만든 향로로 본다",
            "백제의 금속 공예 기술과 정신세계를 보여 주는 대표 유물",
            "국립부여박물관에 소장·전시",
            "‘왕의 무덤에서 나온 순금 유물’ → 능산리 절터 출토, 청동에 금을 입힌 금동",
            "새겨진 인물·동물 무늬가 각각 무엇을 뜻하는지",
            "국립부여박물관",
        ],
    },
    {
        "id": 3, "slug": "cheomseongdae", "heritage": "첨성대", "color": "#345D8C",
        "focus": "첨성대는 정말 천문대였을까?",
        "source": "국가유산청 경주 첨성대 안내",
        "answers": [
            "신라 선덕여왕 때, 7세기로 본다",
            "하늘(천문) 관측과 관련된 시설로 설명된다",
            "남아 있는 가장 오래된 천문 관련 석조 건축물",
            "경주에 원래 자리 그대로 남아 국가유산으로 관리",
            "‘망원경으로 별을 보았다’ → 망원경은 훨씬 뒤에 나왔고 돌을 쌓은 구조물이다",
            "누가 어떤 자세·도구로 관측했는지, 돌 개수의 상징",
            "국가유산청 국가유산포털",
        ],
    },
    {
        "id": 4, "slug": "silla-gold-crown", "heritage": "신라 금관", "color": "#B18A2E",
        "focus": "신라 금관은 실제로 쓰고 다녔을까?",
        "source": "국립중앙박물관 신라 금관 자료",
        "answers": [
            "신라, 5~6세기 무렵 왕릉급 무덤",
            "무덤 주인의 권위를 드러내는 껴묻거리로 만들었다",
            "얇은 금판·세움 장식·굽은옥으로 신라의 황금 문화를 보여 준다",
            "국립중앙박물관·국립경주박물관 등에 소장·전시",
            "‘왕이 평소 매일 썼다’ → 일상 착용 근거가 없고 부장품 성격이 강하다",
            "생전에 실제로 썼는지, 장례를 위해 만들었는지",
            "국립중앙박물관",
        ],
    },
    {
        "id": 5, "slug": "goguryeo-murals", "heritage": "고구려 고분벽화", "color": "#8A4B3F",
        "focus": "고구려 벽화는 무엇을 그린 것일까?",
        "source": "유네스코 고구려 고분군·국립박물관 자료",
        "answers": [
            "고구려, 4~7세기 무덤에 그렸다",
            "왕족·귀족의 무덤 안에 생활과 믿음을 담아 그렸다",
            "기록이 적은 고구려의 생활·복식·믿음을 보여 주는 세계유산",
            "북한과 중국 동북 지역에 나뉘어 있어 사진·모사도로 확인",
            "‘모든 고구려 무덤에 벽화가 있다’ → 일부 무덤에만 남아 있다",
            "벽화 장면이 모든 고구려 사람의 생활인지",
            "유네스코 세계유산센터",
        ],
    },
    {
        "id": 6, "slug": "gaya-tumuli", "heritage": "가야 고분군", "color": "#C87332",
        "focus": "가야는 왜 오랫동안 잊혔을까?",
        "source": "국가유산청 가야고분군·유네스코 자료",
        "answers": [
            "1~6세기 무렵 가야 여러 정치체의 무덤",
            "각 지역 지배층의 무덤으로 만들어졌다",
            "여러 정치체가 나란히 있던 가야를 보여 주는 세계유산(2023년 등재)",
            "김해·함안·고령·합천 등 일곱 고분군으로 남아 관리",
            "‘처음부터 하나의 중앙집권 국가였다’ → 여러 정치체가 함께한 연맹적 성격",
            "고분마다 묻힌 사람의 이름과 일생",
            "국가유산청·유네스코 세계유산센터",
        ],
    },
]


def args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output-root", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--font-dir", type=Path, default=DEFAULT_FONTS)
    return parser.parse_args()


def family(path: Path) -> str:
    font = FontToolsTTFont(str(path), lazy=True)
    try:
        for name_id in (16, 1):
            for record in font["name"].names:
                if record.nameID == name_id:
                    value = record.toUnicode().strip()
                    if value:
                        return value
    finally:
        font.close()
    return path.stem


def setup_fonts(font_dir: Path) -> tuple[str, str]:
    regular = font_dir / "SCDream5.ttf"
    heavy = font_dir / "SCDream9.ttf"
    if not regular.exists() or not heavy.exists():
        regular = Path("/usr/share/fonts/truetype/nanum/NanumGothic.ttf")
        heavy = Path("/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf")
    if not regular.exists() or not heavy.exists():
        raise FileNotFoundError("Korean font files are missing")
    pdfmetrics.registerFont(TTFont("SchoolRegular", str(regular)))
    pdfmetrics.registerFont(TTFont("SchoolHeavy", str(heavy)))
    return family(regular), family(heavy)


def p(text: str, style: ParagraphStyle) -> Paragraph:
    return Paragraph(str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), style)


def styles() -> dict[str, ParagraphStyle]:
    return {
        "body": ParagraphStyle("body", fontName="SchoolRegular", fontSize=7.4, leading=9.4, textColor=TEXT, wordWrap="CJK"),
        "small": ParagraphStyle("small", fontName="SchoolRegular", fontSize=6.6, leading=8.4, textColor=TEXT, wordWrap="CJK"),
        "hint": ParagraphStyle("hint", fontName="SchoolRegular", fontSize=6.0, leading=7.6, textColor=MUTED, wordWrap="CJK"),
        "white": ParagraphStyle("white", fontName="SchoolHeavy", fontSize=7.2, leading=8.6, textColor=WHITE, alignment=1, wordWrap="CJK"),
        "label": ParagraphStyle("label", fontName="SchoolHeavy", fontSize=8.0, leading=9.6, textColor=PRIMARY_DARK, wordWrap="CJK"),
        "answer": ParagraphStyle("answer", fontName="SchoolRegular", fontSize=6.8, leading=8.6, textColor=TEXT, wordWrap="CJK"),
        "center_bold": ParagraphStyle("center_bold", fontName="SchoolHeavy", fontSize=7.0, leading=8.4, textColor=TEXT, alignment=1, wordWrap="CJK"),
    }


def draw_para(c: canvas.Canvas, text: str, style: ParagraphStyle, x: float, y_top: float, w: float, h: float) -> None:
    para = p(text, style)
    _, ph = para.wrap(w, h)
    para.drawOn(c, x, y_top - ph)


def header(c: canvas.Canvas, audience: str) -> float:
    width, height = A4
    x = 13 * mm
    c.setFillColor(PRIMARY)
    c.rect(0, height - 5 * mm, width, 5 * mm, stroke=0, fill=1)
    y = height - 18 * mm
    c.setFillColor(PRIMARY)
    c.roundRect(x, y, 21 * mm, 9 * mm, 3 * mm, stroke=0, fill=1)
    c.setFillColor(WHITE)
    c.setFont("SchoolHeavy", 9.2)
    c.drawCentredString(x + 10.5 * mm, y + 3 * mm, "4차시")
    c.setFillColor(TEXT)
    c.setFont("SchoolHeavy", 17.5)
    c.drawString(x + 27 * mm, y + 1.5 * mm, "우리 모둠 데이터 만들기")
    c.setFont("SchoolRegular", 7.0)
    c.setFillColor(MUTED)
    c.drawRightString(width - x, y + 3 * mm, f"AI혁신학교 공개수업 · {audience}")
    c.setStrokeColor(PRIMARY)
    c.setLineWidth(1.2)
    c.line(x, y - 2 * mm, width - x, y - 2 * mm)
    goal_y = y - 14 * mm
    c.setFillColor(MINT)
    c.rect(x, goal_y, width - 2 * x, 10 * mm, stroke=0, fill=1)
    c.setFillColor(PRIMARY)
    c.rect(x, goal_y, 2.2 * mm, 10 * mm, stroke=0, fill=1)
    c.setFont("SchoolHeavy", 8)
    c.drawString(x + 5 * mm, goal_y + 3.1 * mm, "학습목표")
    c.setFont("SchoolRegular", 7.4)
    c.setFillColor(TEXT)
    c.drawString(x + 26 * mm, goal_y + 3.1 * mm, "여섯 모둠이 같은 일곱 항목으로 조사해 학급 데이터 표의 한 줄을 만들 수 있다.")
    return goal_y - 3.5 * mm


def group_box(c: canvas.Canvas, group: dict[str, Any], y_top: float, st: dict[str, ParagraphStyle]) -> float:
    width, _ = A4
    x = 13 * mm
    w = width - 26 * mm
    h = 22 * mm
    y = y_top - h
    c.setFillColor(WHITE)
    c.setStrokeColor(GRID)
    c.roundRect(x, y, w, h, 2.5 * mm, stroke=1, fill=1)
    accent = HexColor(group["color"])
    c.setFillColor(accent)
    c.roundRect(x, y, 42 * mm, h, 2.5 * mm, stroke=0, fill=1)
    c.setFillColor(WHITE)
    c.setFont("SchoolHeavy", 9)
    c.drawString(x + 4 * mm, y + h - 6.5 * mm, f'{group["id"]}모둠')
    draw_para(c, group["heritage"], ParagraphStyle("heritage", fontName="SchoolHeavy", fontSize=12, leading=14, textColor=WHITE, wordWrap="CJK"), x + 4 * mm, y + h - 9 * mm, 34 * mm, 12 * mm)
    c.setFont("SchoolRegular", 5.8)
    c.drawString(x + 4 * mm, y + 2.6 * mm, ERA["course"])
    rx = x + 46 * mm
    available = w - 46 * mm
    labels = [("핵심 질문", group["focus"]), ("학생", "5학년 ____반  이름 ____________________   역할  □자료  □기록  □기기  □발표")]
    row_h = 11 * mm
    for i, (label, value) in enumerate(labels):
        ry = y + h - (i + 1) * row_h
        c.setFillColor(MINT_LIGHT if i == 0 else WHITE)
        c.rect(rx, ry, available, row_h, stroke=0, fill=1)
        if i:
            c.setStrokeColor(GRID)
            c.line(rx, ry + row_h, x + w - 2 * mm, ry + row_h)
        c.setFont("SchoolHeavy", 6.8)
        c.setFillColor(PRIMARY_DARK)
        c.drawString(rx + 3 * mm, ry + 4 * mm, label)
        draw_para(c, value, st["small"], rx + 20 * mm, ry + row_h - 2.8 * mm, available - 23 * mm, row_h - 1 * mm)
    return y - 4 * mm


def section(c: canvas.Canvas, number: str, title: str, y_top: float) -> float:
    x = 13 * mm
    c.setFillColor(PRIMARY)
    c.roundRect(x, y_top - 7 * mm, 8 * mm, 7 * mm, 2 * mm, stroke=0, fill=1)
    c.setFillColor(WHITE)
    c.setFont("SchoolHeavy", 8.2)
    c.drawCentredString(x + 4 * mm, y_top - 4.8 * mm, number)
    c.setFillColor(TEXT)
    c.setFont("SchoolHeavy", 10.2)
    c.drawString(x + 11 * mm, y_top - 5.2 * mm, title)
    return y_top - 9.5 * mm


def student_page(path: Path, group: dict[str, Any]) -> None:
    """모둠별 A4 한 장 조사 카드. 다른 모둠을 적는 칸은 넣지 않습니다."""
    st = styles()
    path.parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(str(path), pagesize=A4)
    c.setTitle(f'{ERA["short"]} 4차시 {group["id"]}모둠 {group["heritage"]} 조사 카드')
    c.setAuthor("MOAKIT")
    width, _ = A4
    x = 13 * mm
    w = width - 26 * mm
    y = header(c, "학생용 A4 1쪽")
    y = group_box(c, group, y, st)
    y = section(c, "1", "우리 모둠 유산의 일곱 항목 채우기", y)
    c.setFillColor(MUTED)
    c.setFont("SchoolRegular", 6.0)
    c.drawString(x + 11 * mm, y + 1 * mm, "여섯 모둠이 모두 같은 항목을 채워야 학급 데이터 표가 됩니다. 문장 대신 핵심 낱말만 적으세요.")
    y -= 2.5 * mm

    rows = [[p(h, st["white"]) for h in ["조사 항목", "무엇을 찾나요", "핵심 낱말만 적기", "확인한 출처"]]]
    for label, hint in FIELDS:
        rows.append([p(label, st["label"]), p(hint, st["hint"]), p("", st["small"]), p("", st["small"])])
    heights = [9 * mm] + [24.5 * mm] * len(FIELDS)
    table = Table(rows, colWidths=[30 * mm, 38 * mm, w - 108 * mm, 40 * mm], rowHeights=heights)
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), PRIMARY_DARK), ("GRID", (0, 0), (-1, -1), 0.45, GRID),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4), ("TOPPADDING", (0, 0), (-1, -1), 2),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
        ("BACKGROUND", (0, 1), (1, -1), MINT_LIGHT),
        ("BACKGROUND", (0, 7), (-1, 7), BEIGE),
    ]))
    th = sum(heights)
    table.wrapOn(c, w, th)
    table.drawOn(c, x, y - th)
    y -= th + 3 * mm

    c.setFillColor(BEIGE)
    c.roundRect(x, y - 16 * mm, w, 16 * mm, 2 * mm, stroke=0, fill=1)
    c.setFillColor(PRIMARY_DARK)
    c.setFont("SchoolHeavy", 8)
    c.drawString(x + 4 * mm, y - 6.5 * mm, "우리 모둠 카드가 학급 표의 한 줄이 됩니다")
    c.setFillColor(TEXT)
    c.setFont("SchoolRegular", 6.8)
    c.drawString(x + 4 * mm, y - 11.5 * mm, "확인하지 못한 항목은 비워 두거나 ‘아직 모름’이라고 적습니다. 추측으로 채우지 않습니다.")
    c.setFillColor(MUTED)
    c.setFont("SchoolRegular", 5.5)
    c.drawRightString(x + w - 4 * mm, y - 11.5 * mm, "다음 차시: 여섯 모둠 데이터를 한 표로 모아 정제하기")

    c.setFont("SchoolRegular", 5.5)
    c.setFillColor(MUTED)
    c.drawString(x, 7 * mm, "같은 항목으로 모아야 여섯 유산을 비교할 수 있습니다.")
    c.drawRightString(width - x, 7 * mm, f'{ERA["short"]} · {group["id"]}모둠 · 학생용 · S-Core Dream')
    c.showPage()
    c.save()


def merge(paths: Iterable[Path], output: Path) -> None:
    writer = PdfWriter()
    for path in paths:
        for page in PdfReader(str(path)).pages:
            writer.add_page(page)
    with output.open("wb") as f:
        writer.write(f)


def teacher_pdf(path: Path, answers_only: bool = False) -> None:
    st = styles()
    c = canvas.Canvas(str(path), pagesize=A4)
    c.setAuthor("MOAKIT")
    width, _ = A4
    x = 13 * mm
    w = width - 26 * mm

    if not answers_only:
        y = header(c, "교사용 운영안")
        c.setFillColor(TEXT); c.setFont("SchoolHeavy", 17); c.drawString(x, y - 7 * mm, f'{ERA["short"]} 4차시 운영안')
        c.setFillColor(MUTED); c.setFont("SchoolRegular", 7.5)
        c.drawString(x, y - 14 * mm, "PPT·웹앱·활동지가 같은 일곱 항목 이름을 사용합니다. 조사가 아니라 ‘같은 항목으로 데이터 모으기’ 수업입니다.")
        rows = [[p(h, st["white"]) for h in ["시간", "교수·학습 활동", "교사 포인트"]]]
        schedule = [
            ("5분", "지난 시간 우리가 한 확인에 이름 붙이기(출처·시기·교차·원본·보류)", "2차시에 이미 한 일에 이름만 붙입니다."),
            ("5분", "오늘의 미션 안내: 자유 조사가 아니라 같은 항목으로 모으기", "모둠마다 다르게 적으면 비교표가 되지 않음을 보입니다."),
            ("20분", "모둠별로 일곱 항목을 조사하고 웹앱 학급 표에 올리기", "핵심 낱말만 적게 하고 빈칸은 ‘아직 모름’으로 둡니다."),
            ("7분", "모둠 발표: 일곱 항목을 순서대로 읽기", "발표한 모둠부터 정답 슬라이드를 공개합니다."),
            ("3분", "학급 표 확인과 다음 차시 예고", "표를 CSV로 내보내 5차시 시작 파일로 씁니다."),
        ]
        for row in schedule:
            rows.append([p(row[0], st["center_bold"]), p(row[1], st["body"]), p(row[2], st["small"])])
        heights = [9 * mm] + [20 * mm] * len(schedule)
        t = Table(rows, colWidths=[18 * mm, 98 * mm, w - 116 * mm], rowHeights=heights)
        t.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), PRIMARY_DARK), ("GRID", (0, 0), (-1, -1), 0.5, GRID), ("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("LEFTPADDING", (0, 0), (-1, -1), 5), ("RIGHTPADDING", (0, 0), (-1, -1), 5), ("BACKGROUND", (0, 2), (-1, 2), MINT_LIGHT), ("BACKGROUND", (0, 4), (-1, 4), MINT_LIGHT)]))
        th = sum(heights); t.wrapOn(c, w, th); t.drawOn(c, x, y - 21 * mm - th)
        ny = y - 29 * mm - th
        c.setFillColor(BEIGE); c.roundRect(x, ny - 40 * mm, w, 40 * mm, 2 * mm, stroke=0, fill=1)
        c.setFillColor(PRIMARY_DARK); c.setFont("SchoolHeavy", 8); c.drawString(x + 4 * mm, ny - 7 * mm, "학교 양식 적용 기준")
        bullets = [
            "학생 활동지는 모둠별 A4 한 장이며 다른 모둠을 적는 칸이 없습니다.",
            f'일곱 항목 이름은 {" · ".join(label for label, _ in FIELDS)} 로 고정합니다.',
            "학생 화면에는 정답·모범 답안을 표시하지 않고 발표 뒤에 PPT로 공개합니다.",
            "웹앱의 학급 표를 CSV로 내보내 5차시 정제 수업의 시작 파일로 사용합니다.",
        ]
        for i, text in enumerate(bullets):
            c.setFillColor(TEXT); c.setFont("SchoolRegular", 7); c.drawString(x + 5 * mm, ny - (14 + i * 6) * mm, f"• {text}")
        c.setFillColor(MUTED); c.setFont("SchoolRegular", 5.5); c.drawRightString(width - x, 7 * mm, f'{ERA["short"]} · 교사용 · S-Core Dream')
        c.showPage()

    for group in GROUPS:
        y = header(c, "교사용 정답")
        c.setFillColor(HexColor(group["color"])); c.roundRect(x, y - 15 * mm, 28 * mm, 11 * mm, 3 * mm, stroke=0, fill=1)
        c.setFillColor(WHITE); c.setFont("SchoolHeavy", 9); c.drawCentredString(x + 14 * mm, y - 11.5 * mm, f'{group["id"]}모둠')
        c.setFillColor(TEXT); c.setFont("SchoolHeavy", 15); c.drawString(x + 34 * mm, y - 12 * mm, group["heritage"])
        c.setFont("SchoolRegular", 7); c.setFillColor(MUTED); c.drawRightString(width - x, y - 11.5 * mm, group["focus"])
        rows = [[p(h, st["white"]) for h in ["조사 항목", "모범 답안 (핵심 낱말)"]]]
        for (label, _), answer in zip(FIELDS, group["answers"]):
            rows.append([p(label, st["label"]), p(answer, st["answer"])])
        heights = [9 * mm] + [22 * mm] * len(FIELDS)
        t = Table(rows, colWidths=[34 * mm, w - 34 * mm], rowHeights=heights)
        t.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), PRIMARY_DARK), ("GRID", (0, 0), (-1, -1), 0.45, GRID), ("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("LEFTPADDING", (0, 0), (-1, -1), 5), ("RIGHTPADDING", (0, 0), (-1, -1), 5), ("BACKGROUND", (0, 1), (0, -1), MINT_LIGHT), ("BACKGROUND", (0, 7), (-1, 7), BEIGE)]))
        th = sum(heights); t.wrapOn(c, w, th); t.drawOn(c, x, y - 20 * mm - th)
        sy = y - 25 * mm - th
        c.setFillColor(MINT); c.roundRect(x, sy - 14 * mm, w, 14 * mm, 2 * mm, stroke=0, fill=1)
        c.setFillColor(PRIMARY_DARK); c.setFont("SchoolHeavy", 7.3); c.drawString(x + 4 * mm, sy - 6 * mm, "우선 확인처")
        c.setFillColor(TEXT); c.setFont("SchoolRegular", 6.8); c.drawString(x + 28 * mm, sy - 6 * mm, group["source"])
        c.setFillColor(MUTED); c.setFont("SchoolRegular", 5.8); c.drawString(x + 4 * mm, sy - 11 * mm, "학생 발표를 들은 뒤에 공개하고, 학생이 적은 낱말과 범위를 비교합니다.")
        c.setFillColor(MUTED); c.setFont("SchoolRegular", 5.5); c.drawRightString(width - x, 7 * mm, f'{ERA["short"]} · {group["id"]}모둠 · 정답 · S-Core Dream')
        c.showPage()
    c.save()


def add_text(slide, x: float, y: float, w: float, h: float, text: str, font: str, size: float, color: RGBColor = PPT_TEXT, bold: bool = False, align=PP_ALIGN.LEFT, fill: RGBColor | None = None) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame; tf.clear()
    tf.margin_left = Inches(0.07); tf.margin_right = Inches(0.07); tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]; para.alignment = align; para.text = text
    for run in para.runs:
        run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color


def rect(slide, x: float, y: float, w: float, h: float, fill: RGBColor, rounded: bool = False) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()


def ppt_header(slide, regular: str, heavy: str) -> None:
    rect(slide, 0, 0, 13.333, 0.12, PPT_PRIMARY)
    rect(slide, 0.55, 0.36, 1.15, 0.42, PPT_PRIMARY, True)
    add_text(slide, 0.55, 0.36, 1.15, 0.42, "4차시", heavy, 12, PPT_WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 1.88, 0.28, 6.6, 0.58, "우리 모둠 데이터 만들기", heavy, 25, PPT_TEXT, True)
    add_text(slide, 9.4, 0.35, 3.35, 0.35, f'{ERA["short"]} · 교사용 수업 PPT', regular, 9, PPT_MUTED, False, PP_ALIGN.RIGHT)
    rect(slide, 0.55, 0.93, 12.23, 0.025, PPT_PRIMARY)


def ppt(path: Path, regular: str, heavy: str) -> None:
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(blank); ppt_header(slide, regular, heavy)
        add_text(slide, 0.72, 1.18, 11.9, 0.62, title, heavy, 28, PPT_TEXT, True)
        if subtitle:
            add_text(slide, 0.75, 1.77, 11.8, 0.38, subtitle, regular, 12, PPT_MUTED)
        return slide

    # 1 · 표지
    slide = prs.slides.add_slide(blank); rect(slide, 0, 0, 13.333, 7.5, RGBColor(248,251,249)); rect(slide, 0, 0, 13.333, 0.2, PPT_PRIMARY)
    add_text(slide, 0.82, 0.75, 6, 0.36, "AI혁신학교 공개수업 · 4차시", regular, 12, PPT_DARK, True)
    add_text(slide, 0.82, 1.38, 10.8, 1.18, "우리 모둠 데이터 만들기", heavy, 42, PPT_TEXT, True)
    rect(slide, 0.82, 3.05, 11.55, 1.38, PPT_MINT); rect(slide, 0.82, 3.05, 0.1, 1.38, PPT_PRIMARY)
    add_text(slide, 1.18, 3.28, 10.7, 0.35, "핵심 질문", heavy, 13, PPT_DARK, True)
    add_text(slide, 1.18, 3.68, 10.7, 0.55, f'“{ERA["core"]}”', heavy, 22, PPT_DARK, True)
    add_text(slide, 0.82, 5.05, 11.5, 0.95, "여섯 모둠이 똑같은 일곱 항목을 채우면\n우리 모둠 카드가 학급 데이터 표의 한 줄이 됩니다.", regular, 21)

    # 2 · 도입 (지난 시간 확인에 이름 붙이기 · 5분)
    slide = new("지난 시간 우리가 한 확인에 이름 붙이기", "2차시에 이미 해 본 일입니다. 이름만 붙여 봅시다.")
    for i, (name, detail) in enumerate(VERIFICATION_NAMES):
        x = 0.85 + i * 2.42
        rect(slide, x, 2.3, 2.15, 1.05, PPT_PRIMARY, True)
        add_text(slide, x, 2.3, 2.15, 1.05, name, heavy, 24, PPT_WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, x, 3.45, 2.15, 1.15, detail, regular, 13, PPT_TEXT, False, PP_ALIGN.CENTER)
    rect(slide, 0.85, 5.0, 11.6, 1.05, PPT_BEIGE, True)
    add_text(slide, 1.1, 5.0, 11.1, 1.05, "이 다섯 가지로 확인한 내용을 오늘은 ‘같은 항목’에 나누어 담습니다.", heavy, 19, PPT_DARK, True, PP_ALIGN.CENTER)

    # 3 · 오늘의 미션
    slide = new("오늘의 미션", "조사하기가 아니라 데이터 만들기입니다.")
    rect(slide, 0.85, 2.2, 5.65, 2.5, RGBColor(252,239,239), True)
    add_text(slide, 1.05, 2.4, 5.25, 0.45, "자유롭게 조사하면", heavy, 18, RGBColor(154,62,62), True)
    add_text(slide, 1.05, 2.95, 5.25, 1.6, "모둠마다 다른 내용을 적어서\n여섯 유산을 나란히 놓고\n비교할 수 없습니다.", regular, 17, PPT_TEXT)
    rect(slide, 6.8, 2.2, 5.65, 2.5, PPT_MINT, True)
    add_text(slide, 7.0, 2.4, 5.25, 0.45, "같은 항목으로 모으면", heavy, 18, PPT_DARK, True)
    add_text(slide, 7.0, 2.95, 5.25, 1.6, "여섯 줄이 하나의 표가 되어\n시기·가치·현재 상태를\n서로 비교할 수 있습니다.", regular, 17, PPT_TEXT)
    rect(slide, 0.85, 5.05, 11.6, 1.1, PPT_PRIMARY, True)
    add_text(slide, 1.1, 5.05, 11.1, 1.1, "우리 모둠 카드가 학급 데이터 표의 한 줄이 됩니다", heavy, 24, PPT_WHITE, True, PP_ALIGN.CENTER)

    # 4 · 일곱 항목 (활동지와 같은 표 모양)
    slide = new("우리가 채울 일곱 항목", "활동지 표와 같은 순서, 같은 이름입니다.")
    rect(slide, 0.85, 2.12, 2.6, 0.42, PPT_DARK)
    add_text(slide, 0.85, 2.12, 2.6, 0.42, "조사 항목", heavy, 12, PPT_WHITE, True)
    rect(slide, 3.5, 2.12, 8.95, 0.42, PPT_DARK)
    add_text(slide, 3.7, 2.12, 8.75, 0.42, "무엇을 찾나요", heavy, 12, PPT_WHITE, True)
    for i, (label, hint) in enumerate(FIELDS):
        y = 2.6 + i * 0.62
        band = PPT_BEIGE if label == "출처" else (PPT_MINT if i % 2 == 0 else RGBColor(247,249,248))
        rect(slide, 0.85, y, 2.6, 0.55, band)
        add_text(slide, 0.85, y, 2.6, 0.55, label, heavy, 15, PPT_DARK, True, PP_ALIGN.CENTER)
        rect(slide, 3.5, y, 8.95, 0.55, RGBColor(252,252,251))
        add_text(slide, 3.7, y, 8.75, 0.55, hint, regular, 14, PPT_TEXT)

    # 5 · 조사 시간 안내
    slide = new("조사 시간 안내 · 20분", "공식 자료부터 순서대로 확인합니다.")
    for i, text in enumerate(SOURCE_RANKS):
        rect(slide, 0.85, 2.2 + i * 0.78, 6.9, 0.62, PPT_PRIMARY if i < 3 else PPT_BEIGE, True)
        add_text(slide, 1.1, 2.2 + i * 0.78, 6.4, 0.62, text, heavy if i < 3 else regular, 15, PPT_WHITE if i < 3 else PPT_TEXT, i < 3)
    rect(slide, 8.1, 2.2, 4.35, 2.96, PPT_MINT, True)
    add_text(slide, 8.35, 2.35, 3.85, 0.4, "이렇게 적습니다", heavy, 15, PPT_DARK, True)
    add_text(slide, 8.35, 2.8, 3.85, 2.2, "• 문장 말고 핵심 낱말만\n• 항목마다 확인한 출처도 함께\n• 확인 못 한 항목은\n   ‘아직 모름’으로 남기기", regular, 14, PPT_TEXT)
    rect(slide, 0.85, 5.45, 11.6, 1.0, RGBColor(247,249,248), True)
    add_text(slide, 1.1, 5.45, 11.1, 1.0, f'웹앱: 수업 사이트 → {ERA["short"]} → 4차시 → ‘{ERA["appTab"]}’ 탭 (로그인 없이 모둠만 선택)', heavy, 17, PPT_DARK, True, PP_ALIGN.CENTER)

    # 6 · 발표 안내
    slide = new("모둠 발표 안내", "일곱 항목을 순서대로 읽으면 그대로 발표가 됩니다.")
    rect(slide, 0.85, 2.2, 11.6, 2.9, PPT_MINT, True)
    add_text(slide, 1.2, 2.45, 10.9, 0.5, "우리 모둠은 ___모둠, ___을 조사했습니다.", regular, 19)
    add_text(slide, 1.2, 3.0, 10.9, 0.5, "시기는 ___, 만든 까닭은 ___, 가치는 ___입니다.", regular, 19)
    add_text(slide, 1.2, 3.55, 10.9, 0.5, "현재 상태는 ___이고, AI가 틀린 부분은 ___로 고쳤습니다.", regular, 19)
    add_text(slide, 1.2, 4.1, 10.9, 0.6, "아직 모르는 점은 ___이고, 출처는 ___입니다.", heavy, 20, PPT_DARK, True)
    rect(slide, 0.85, 5.4, 11.6, 0.95, PPT_BEIGE, True)
    add_text(slide, 1.1, 5.4, 11.1, 0.95, "발표를 들은 뒤에 그 모둠의 정답 슬라이드를 엽니다.", heavy, 18, PPT_DARK, True, PP_ALIGN.CENTER)

    # 7~12 · 모둠별 정답 (발표 뒤에만 공개)
    for group in GROUPS:
        slide = new(f'{group["id"]}모둠 정답 · {group["heritage"]}', "이 모둠의 발표를 들은 뒤에 공개하세요.")
        rect(slide, 10.75, 0.3, 1.95, 0.52, RGBColor(161,73,73), True)
        add_text(slide, 10.75, 0.3, 1.95, 0.52, "발표 후 공개", heavy, 12, PPT_WHITE, True, PP_ALIGN.CENTER)
        color = RGBColor(*tuple(int(group["color"][j:j+2], 16) for j in (1,3,5)))
        rect(slide, 0.82, 2.2, 12.05, 0.02, color)
        for i, ((label, _), answer) in enumerate(zip(FIELDS, group["answers"])):
            y = 2.32 + i * 0.63
            rect(slide, 0.82, y, 2.35, 0.55, color if label == "출처" else PPT_PRIMARY, True)
            add_text(slide, 0.82, y, 2.35, 0.55, label, heavy, 13, PPT_WHITE, True, PP_ALIGN.CENTER)
            add_text(slide, 3.35, y, 9.5, 0.55, answer, regular, 13.5, PPT_TEXT)
        add_text(slide, 0.82, 6.9, 12.05, 0.36, f'우선 확인처: {group["source"]}', regular, 11, PPT_MUTED)

    # 13 · 정리 · 다음 차시 예고
    slide = new("여섯 모둠 데이터를 한 표로 모으면?", "다음 시간에 학급 데이터 표를 정제합니다.")
    header_cells = ["모둠", "유산"] + [label for label, _ in FIELDS]
    widths = [0.95, 1.75, 1.25, 1.45, 1.15, 1.45, 1.85, 1.25, 1.05]
    left = 0.72
    for cell, cell_width in zip(header_cells, widths):
        rect(slide, left, 2.25, cell_width - 0.05, 0.5, PPT_DARK)
        add_text(slide, left, 2.25, cell_width - 0.05, 0.5, cell, heavy, 10.5, PPT_WHITE, True, PP_ALIGN.CENTER)
        left += cell_width
    for i, group in enumerate(GROUPS):
        y = 2.8 + i * 0.44
        left = 0.72
        for j, cell_width in enumerate(widths):
            rect(slide, left, y, cell_width - 0.05, 0.4, PPT_MINT if i % 2 == 0 else RGBColor(250,251,250))
            if j == 0:
                add_text(slide, left, y, cell_width - 0.05, 0.4, f'{group["id"]}모둠', heavy, 10, PPT_DARK, True, PP_ALIGN.CENTER)
            elif j == 1:
                add_text(slide, left, y, cell_width - 0.05, 0.4, group["heritage"], regular, 9.5, PPT_TEXT, False, PP_ALIGN.CENTER)
            left += cell_width
    rect(slide, 0.72, 5.6, 11.9, 1.05, PPT_BEIGE, True)
    add_text(slide, 1.0, 5.6, 11.4, 1.05, "다음 시간 · 표기가 다른 칸, 빈칸, 겹치는 줄을 찾아 하나의 기준으로 정제합니다.", heavy, 19, PPT_DARK, True, PP_ALIGN.CENTER)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def rebuild(zip_path: Path, files: Iterable[Path], base: Path) -> None:
    temp = zip_path.with_suffix(".tmp.zip")
    with zipfile.ZipFile(temp, "w", zipfile.ZIP_DEFLATED, compresslevel=9) as archive:
        for file in sorted({f for f in files if f.exists()}):
            if file.resolve() in {zip_path.resolve(), temp.resolve()}:
                continue
            archive.write(file, file.relative_to(base))
    temp.replace(zip_path)


def update_manifest(output: Path) -> None:
    path = output / "manifest.json"
    if not path.exists():
        return
    data = json.loads(path.read_text(encoding="utf-8"))
    era_record = data.get("eras", {}).get(ERA["id"])
    if not era_record:
        return
    lesson = next((item for item in era_record.get("lessons", []) if item.get("lessonId") == 4), None)
    if not lesson:
        return
    files = lesson.setdefault("files", {})
    for key, name in [("student", "lesson-04-student.pdf"), ("teacher", "lesson-04-teacher.pdf"), ("answer", "lesson-04-answer.pdf"), ("bundle", "lesson-04-all.zip"), ("ppt", "lesson-04-teaching.pptx")]:
        file = output / ERA["id"] / name
        if file.exists():
            files[key] = {"path": f'/downloads/{ERA["id"]}/{name}', "size": file.stat().st_size}
    files["groups"] = [
        {
            "groupId": group["id"],
            "heritage": group["heritage"],
            "path": f'/downloads/{ERA["id"]}/lesson-04-group-{group["id"]:02d}-{group["slug"]}.pdf',
            "size": (output / ERA["id"] / f'lesson-04-group-{group["id"]:02d}-{group["slug"]}.pdf').stat().st_size,
            "pages": 1,
            "format": "A4",
        }
        for group in GROUPS
    ]
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def main() -> None:
    options = args()
    output = options.output_root
    regular_name, _heavy_name = setup_fonts(options.font_dir)
    heavy_name = _heavy_name
    directory = output / ERA["id"]
    directory.mkdir(parents=True, exist_ok=True)

    groups: list[Path] = []
    for group in GROUPS:
        file = directory / f'lesson-04-group-{group["id"]:02d}-{group["slug"]}.pdf'
        student_page(file, group)
        groups.append(file)

    student = directory / "lesson-04-student.pdf"; merge(groups, student)
    teacher = directory / "lesson-04-teacher.pdf"; teacher_pdf(teacher, False)
    answer = directory / "lesson-04-answer.pdf"; teacher_pdf(answer, True)
    deck = directory / "lesson-04-teaching.pptx"; ppt(deck, regular_name, heavy_name)
    rebuild(directory / "lesson-04-all.zip", [student, teacher, answer, deck, *groups], directory)
    era_zip = directory / f'{ERA["id"]}-all-materials.zip'
    rebuild(era_zip, [f for f in directory.iterdir() if f.is_file() and f.name != era_zip.name], directory)
    update_manifest(output)
    print(f"Generated lesson 4 data-card materials in {directory}")


if __name__ == "__main__":
    main()
